from pptx import Presentation
from pptx.util import Cm

def create_presentation():
    """
    Create a new presentation with A4 size
    """
    p = Presentation()
    p.slide_width = Cm(21)
    p.slide_height = Cm(29.7)
    return p

def add_slide(presentation):
    """
    Add a blank slide to the presentation
    """
    return presentation.slides.add_slide(presentation.slide_layouts[6])