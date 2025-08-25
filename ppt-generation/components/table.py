from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from config.colors import TABLE_HEADER_BG, TABLE_HEADER_TEXT, TABLE_ALT_ROW, TEXT_DARK
from config.sizes import TABLE_LEFT, TABLE_TOP, TABLE_WIDTH, TABLE_HEIGHT, TABLE_COL_WIDTHS
from config.sizes import SECOND_TABLE_LEFT, SECOND_TABLE_TOP, SECOND_TABLE_WIDTH, SECOND_TABLE_HEIGHT, SECOND_TABLE_COL_WIDTHS
from config.sizes import TABLE_HEADER_FONT, TABLE_BODY_FONT
from pptx.util import Cm


def create_table(slide, table_data, rows, cols, 
                left=TABLE_LEFT, top=TABLE_TOP, 
                width=TABLE_WIDTH, height=TABLE_HEIGHT, 
                column_widths=TABLE_COL_WIDTHS, is_second_table=False):
    """
    Create a table on the slide
    """
    table_shape = slide.shapes.add_table(
        rows=rows,
        cols=cols,
        left=left,
        top=top,
        width=width,
        height=height
    )

    table = table_shape.table

    # Set column widths
    for i, width in enumerate(column_widths):
        table.columns[i].width = width

    # Populate table and style
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_data
            
            # Style header row
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
                
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                paragraph = cell.text_frame.paragraphs[0]
                run = paragraph.runs[0]
                run.font.color.rgb = TABLE_HEADER_TEXT
                run.font.bold = True
                run.font.size = TABLE_HEADER_FONT
                paragraph.alignment = PP_ALIGN.LEFT
            else:
                # Data row styling
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TABLE_ALT_ROW
                
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                paragraph = cell.text_frame.paragraphs[0]
                run = paragraph.runs[0]
                run.font.size = TABLE_BODY_FONT
                run.font.color.rgb = TEXT_DARK
                
                # For second table, make first column bold
                if is_second_table and col_idx == 0:
                    run.font.bold = True
                
                paragraph.alignment = PP_ALIGN.LEFT

    # Add table borders
    for row in table.rows:
        for cell in row.cells:
            cell.margin_left = Cm(0.1)
            cell.margin_right = Cm(0.1)
            cell.margin_top = Cm(0.05)
            cell.margin_bottom = Cm(0.05)

    return table