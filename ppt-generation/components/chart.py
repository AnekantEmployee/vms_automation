from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from config.colors import CHART_COLORS, GRAY, WHITE, TEXT_DARK
from config.sizes import CHART_LEFT, CHART_TOP, CHART_WIDTH, CHART_HEIGHT, CHART_TITLE_FONT, CHART_AXIS_FONT, CHART_LABEL_FONT

def create_chart(slide, chart_data, title="Vulnerability Count", 
                left=CHART_LEFT, top=CHART_TOP, 
                width=CHART_WIDTH, height=CHART_HEIGHT):
    """
    Create a chart on the slide
    """
    # Create chart data
    chart_data_obj = CategoryChartData()
    chart_data_obj.categories = chart_data['categories']
    chart_data_obj.add_series('Vulnerabilities', chart_data['values'])

    # Add chart
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left, top, width, height,
        chart_data_obj
    )

    # Customize chart
    chart_obj = chart.chart

    # Set chart title
    chart_obj.has_title = True
    chart_obj.chart_title.text_frame.text = title
    chart_obj.chart_title.text_frame.paragraphs[0].font.size = CHART_TITLE_FONT
    chart_obj.chart_title.text_frame.paragraphs[0].font.color.rgb = GRAY

    # Get series and customize colors
    series = chart_obj.series[0]

    try:
        points = series.points
        for i, point in enumerate(points):
            if i < len(CHART_COLORS):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = CHART_COLORS[i]
    except Exception:
        # Fallback if individual point coloring fails
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = CHART_COLORS[2]  # Medium orange

    # Customize axes
    try:
        value_axis = chart_obj.value_axis
        value_axis.maximum_scale = 5.5
        value_axis.minimum_scale = 0
        value_axis.major_unit = 1
        
        # Clean up axis appearance
        value_axis.tick_labels.font.size = CHART_AXIS_FONT
        category_axis = chart_obj.category_axis
        category_axis.tick_labels.font.size = CHART_AXIS_FONT
        
        # Make gridlines white
        value_axis.major_gridlines.format.line.color.rgb = WHITE
    except Exception:
        pass

    # Add data labels
    try:
        series.has_data_labels = True
        data_labels = series.data_labels
        data_labels.show_value = True
        data_labels.font.size = CHART_LABEL_FONT
        data_labels.font.color.rgb = TEXT_DARK
    except Exception:
        pass

    # Remove legend and clean up appearance
    try:
        chart_obj.has_legend = False
        chart_obj.plot_area.format.line.color.rgb = WHITE
        chart.format.line.color.rgb = WHITE
    except Exception:
        pass

    return chart