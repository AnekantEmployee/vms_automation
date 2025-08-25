from pptx.util import Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from config.colors import HEADER_BG, WHITE, RED, TEXT_BLACK
from config.sizes import FOOTER_FULL_LEFT, FOOTER_FULL_TOP, FOOTER_FULL_WIDTH, FOOTER_FULL_HEIGHT
from config.sizes import DISCLAIMER_LEFT, DISCLAIMER_TOP, DISCLAIMER_WIDTH, DISCLAIMER_HEIGHT, DISCLAIMER_FONT
from config.sizes import TAGLINE_LEFT, TAGLINE_TOP, TAGLINE_WIDTH, TAGLINE_HEIGHT, TAGLINE_FONT
from config.sizes import WHAT_IN_IT_LEFT, WHAT_IN_IT_TOP, WHAT_IN_IT_WIDTH, WHAT_IN_IT_HEIGHT, BULLET_FONT

# Define blue color for icons
ICON_BLUE = RGBColor(0, 162, 232)  # Bright blue color

def create_yash_recommendations_section(slide, RECOMMENDATIONS_DATA):
    """
    Create the "YASH Recommendations - Why Learn More?" section (40% width on right side)
    with proper blue circular icons
    """
    # Fixed positioning to match the image layout
    recommendations_left = Cm(14.5)  # Position on right side
    recommendations_top = Cm(4.5)   # Align with content area
    recommendations_width = Cm(9.5)  # 40% of slide width
    recommendations_height = Cm(10)  # Appropriate height
    
    recommendations_box = slide.shapes.add_textbox(
        recommendations_left, recommendations_top, recommendations_width, recommendations_height
    )

    # Style the text box
    recommendations_frame = recommendations_box.text_frame
    recommendations_frame.word_wrap = True
    recommendations_frame.margin_left = Cm(0.5)
    recommendations_frame.margin_right = Cm(0.5)
    recommendations_frame.margin_top = Cm(0.5)
    recommendations_frame.margin_bottom = Cm(0.5)

    # Add main title
    title_para = recommendations_frame.paragraphs[0]
    title_run = title_para.add_run()
    title_run.text = "YASH Recommendations"
    title_run.font.size = Cm(0.7)
    title_run.font.bold = True
    title_run.font.color.rgb = TEXT_BLACK
    title_para.alignment = PP_ALIGN.LEFT
    
    # Add subtitle
    subtitle_para = recommendations_frame.add_paragraph()
    subtitle_run = subtitle_para.add_run()
    subtitle_run.text = "Why Learn More?"
    subtitle_run.font.size = Cm(0.6)
    subtitle_run.font.color.rgb = TEXT_BLACK
    subtitle_para.alignment = PP_ALIGN.LEFT
    
    # Add spacing
    recommendations_frame.add_paragraph()
    
    # Create circular blue backgrounds for icons and text
    current_y = recommendations_top + Cm(2.5)  # Start position for first item
    
    for i, rec in enumerate(RECOMMENDATIONS_DATA):
        # Create blue circle for icon
        circle_left = recommendations_left + Cm(0.2)
        circle_top = current_y
        circle_size = Cm(1.2)
        
        # Add blue circle shape
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            circle_left, circle_top, circle_size, circle_size
        )
        
        # Style the circle
        circle.fill.solid()
        circle.fill.fore_color.rgb = ICON_BLUE
        circle.line.fill.background()  # Remove border
        
        # Add icon text in center of circle
        icon_box = slide.shapes.add_textbox(
            circle_left, circle_top, circle_size, circle_size
        )
        icon_frame = icon_box.text_frame
        icon_frame.margin_left = 0
        icon_frame.margin_right = 0
        icon_frame.margin_top = 0
        icon_frame.margin_bottom = 0
        
        icon_para = icon_frame.paragraphs[0]
        icon_run = icon_para.add_run()
        icon_run.text = rec['icon']
        icon_run.font.size = Cm(0.6)
        icon_para.alignment = PP_ALIGN.CENTER
        
        # Make icon background transparent
        icon_box.fill.background()
        
        # Add title and subtitle text next to circle
        text_left = circle_left + circle_size + Cm(0.5)
        text_top = circle_top
        text_width = Cm(6)
        text_height = Cm(1.5)
        
        text_box = slide.shapes.add_textbox(
            text_left, text_top, text_width, text_height
        )
        
        text_frame = text_box.text_frame
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = Cm(0.1)
        text_frame.margin_bottom = 0
        
        # Add title
        title_para = text_frame.paragraphs[0]
        title_run = title_para.add_run()
        title_run.text = rec['title']
        title_run.font.size = Cm(0.5)
        title_run.font.bold = True
        title_run.font.color.rgb = TEXT_BLACK
        title_para.alignment = PP_ALIGN.LEFT
        
        # Add subtitle
        subtitle_para = text_frame.add_paragraph()
        subtitle_run = subtitle_para.add_run()
        subtitle_run.text = rec['subtitle']
        subtitle_run.font.size = Cm(0.4)
        subtitle_run.font.color.rgb = TEXT_BLACK
        subtitle_para.alignment = PP_ALIGN.LEFT
        
        # Make text background transparent
        text_box.fill.background()
        
        # Move to next position
        current_y += Cm(2)
    
    # Add framework note at the bottom
    framework_left = recommendations_left + Cm(0.5)
    framework_top = current_y + Cm(0.5)
    framework_width = Cm(8)
    framework_height = Cm(1)
    
    framework_box = slide.shapes.add_textbox(
        framework_left, framework_top, framework_width, framework_height
    )
    
    framework_frame = framework_box.text_frame
    framework_para = framework_frame.paragraphs[0]
    framework_run = framework_para.add_run()
    framework_run.text = "(e.g. NIST & ISO/IEC 27001)"
    framework_run.font.size = Cm(0.35)
    framework_run.font.color.rgb = TEXT_BLACK
    framework_run.font.italic = True
    framework_para.alignment = PP_ALIGN.LEFT
    
    # Make background transparent
    recommendations_box.fill.background()
    framework_box.fill.background()
    
    return recommendations_box

def create_whats_in_it_section(slide):
    """
    Create the "What's in it for you?" section (60% width on left side)
    """
    # Fixed positioning to align with the left content area
    whats_in_it_left = WHAT_IN_IT_LEFT  # Use existing config value
    whats_in_it_width = Cm(13)  # Fixed width for left section (60% of slide)
    
    whats_in_it_box = slide.shapes.add_textbox(
        whats_in_it_left, WHAT_IN_IT_TOP, whats_in_it_width, WHAT_IN_IT_HEIGHT
    )

    # Style the text box
    whats_in_it_frame = whats_in_it_box.text_frame
    whats_in_it_frame.word_wrap = True
    whats_in_it_frame.margin_left = Cm(0.5)
    whats_in_it_frame.margin_right = Cm(0.5)
    whats_in_it_frame.margin_top = Cm(0.5)
    whats_in_it_frame.margin_bottom = Cm(0.5)

    # Add title
    title_para = whats_in_it_frame.paragraphs[0]
    title_run = title_para.add_run()
    title_run.text = "What's in it for you?"
    title_run.font.size = Cm(0.6)
    title_run.font.bold = True
    title_run.font.color.rgb = TEXT_BLACK
    title_para.alignment = PP_ALIGN.LEFT
    
    # Add spacing
    whats_in_it_frame.add_paragraph()
    
    # Add bullet points
    bullet_points = [
        "Explore the most prevalent and impactful threats, techniques, and trends that we've observed.",
        "Prioritize and categorize the threats based on their severity and impact on your business and invest your valuable time and resources on stuff that needs immediate attention.",
        "Shape and inform your readiness, detection, and response to critical threats",
        "Get a custom consultation from our cybersecurity experts and secure your IT landscape."
    ]
    
    for bullet_text in bullet_points:
        bullet_para = whats_in_it_frame.add_paragraph()
        bullet_para.level = 0
        bullet_run = bullet_para.add_run()
        bullet_run.text = f"• {bullet_text}"  # Using bullet point character
        bullet_run.font.size = BULLET_FONT
        bullet_run.font.color.rgb = TEXT_BLACK
        bullet_para.alignment = PP_ALIGN.LEFT

    # Add spacing
    for _ in range(3):
        whats_in_it_frame.add_paragraph()

    # Add contact info
    contact_para = whats_in_it_frame.add_paragraph()
    contact_run1 = contact_para.add_run()
    contact_run1.text = "For more information write to us at"
    contact_run1.font.size = BULLET_FONT
    contact_run1.font.color.rgb = TEXT_BLACK
    contact_para.alignment = PP_ALIGN.LEFT

    contact_para2 = whats_in_it_frame.add_paragraph()
    contact_run2 = contact_para2.add_run()
    contact_run2.text = "cybersecurity@yash.com"
    contact_run2.font.size = Cm(0.4)
    contact_run2.font.bold = True
    contact_run2.font.color.rgb = TEXT_BLACK
    contact_para2.alignment = PP_ALIGN.LEFT

    # Make background transparent
    whats_in_it_box.fill.background()
    
    return whats_in_it_box

def create_footer(slide):
    """
    Create the footer section with proper alignment
    """
    # Add full width blue strip
    footer_full = slide.shapes.add_textbox(
        FOOTER_FULL_LEFT, FOOTER_FULL_TOP, FOOTER_FULL_WIDTH, FOOTER_FULL_HEIGHT
    )

    # Add blue background to full strip
    footer_full.fill.solid()
    footer_full.fill.fore_color.rgb = HEADER_BG

    # Add disclaimer text (left 60% to match content layout)
    disclaimer_width = Cm(13)  # Match left section width
    disclaimer_left = slide.shapes.add_textbox(
        DISCLAIMER_LEFT, DISCLAIMER_TOP, disclaimer_width, DISCLAIMER_HEIGHT
    )

    # Style disclaimer text
    disclaimer_frame = disclaimer_left.text_frame
    disclaimer_frame.word_wrap = True
    disclaimer_frame.margin_left = Cm(0.5)
    disclaimer_frame.margin_right = Cm(0.5)
    disclaimer_frame.margin_top = Cm(0.2)
    disclaimer_frame.margin_bottom = Cm(0.2)

    # Add disclaimer text with red "Disclaimer:" 
    disclaimer_para = disclaimer_frame.paragraphs[0]
    disclaimer_run1 = disclaimer_para.add_run()
    disclaimer_run1.text = "Disclaimer:"
    disclaimer_run1.font.size = DISCLAIMER_FONT
    disclaimer_run1.font.color.rgb = RED

    # Add rest of text in white
    disclaimer_run2 = disclaimer_para.add_run()
    disclaimer_run2.text = " All the information for the report has been obtained from publicly available sources. YASH technologies does not perform any unauthorized assessment that could impact your business"
    disclaimer_run2.font.size = DISCLAIMER_FONT
    disclaimer_run2.font.color.rgb = WHITE

    disclaimer_para.alignment = PP_ALIGN.LEFT

    # Add tagline text (right 40% to match recommendations layout)
    tagline_left = Cm(14.5)  # Match recommendations left position
    tagline_width = Cm(9.5)  # Match recommendations width
    tagline_right = slide.shapes.add_textbox(
        tagline_left, TAGLINE_TOP, tagline_width, TAGLINE_HEIGHT
    )

    # Style tagline text
    tagline_frame = tagline_right.text_frame
    tagline_frame.margin_left = Cm(0.5)
    tagline_frame.margin_right = Cm(0.5)
    tagline_frame.margin_top = Cm(0.2)
    tagline_frame.margin_bottom = Cm(0.2)

    tagline_frame.text = "Maximise security posture with"
    tagline_para = tagline_frame.paragraphs[0]
    tagline_run = tagline_para.runs[0]
    tagline_run.font.size = TAGLINE_FONT
    tagline_run.font.color.rgb = WHITE
    tagline_run.font.bold = True
    tagline_para.alignment = PP_ALIGN.CENTER

    # Add second line with white color for SOC/MDR/VMS/TRPM
    tagline_para2 = tagline_frame.add_paragraph()
    tagline_run2 = tagline_para2.add_run()
    tagline_run2.text = "SOC/MDR/VMS/TRPM"
    tagline_run2.font.size = TAGLINE_FONT
    tagline_run2.font.color.rgb = WHITE
    tagline_run2.font.bold = True
    tagline_para2.alignment = PP_ALIGN.CENTER

    # Make background transparent for text boxes
    disclaimer_left.fill.background()
    tagline_right.fill.background()
    
    return footer_full, disclaimer_left, tagline_right

def create_complete_slide_layout(slide, RECOMMENDATIONS_DATA):
    """
    Create the complete slide layout with both sections properly aligned
    """
    # Create left section (60% width) - "What's in it for you?"
    whats_in_it = create_whats_in_it_section(slide)
    
    # Create right section (40% width) - "YASH Recommendations"
    yash_recommendations = create_yash_recommendations_section(slide, RECOMMENDATIONS_DATA)
    
    # Create footer with proper alignment
    footer_elements = create_footer(slide)
    
    return whats_in_it, yash_recommendations, footer_elements
