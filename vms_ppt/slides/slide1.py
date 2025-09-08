from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from .slide_utils import SlideUtils
from config_colors import COLORS, FONT_SIZES

def create_slide1(prs: Presentation, slide1_data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Create title bar using utility
    SlideUtils.create_title_bar(slide, prs, slide1_data["title"])
    
    # Create agenda box
    _create_agenda_box(slide, prs, slide1_data["agenda_points"])

def _create_agenda_box(slide, prs, agenda_points):
    """Create agenda box with points"""
    agenda_left = Inches(1)
    agenda_top = Inches(1.2)
    agenda_width = prs.slide_width - Inches(2)
    agenda_height = Inches(5.5)
    
    agenda_box = slide.shapes.add_textbox(agenda_left, agenda_top, agenda_width, agenda_height)
    agenda_tf = agenda_box.text_frame
    agenda_tf.word_wrap = True
    
    # Style the box
    agenda_box.fill.solid()
    agenda_box.fill.fore_color.rgb = RGBColor(*COLORS["light_gray"])
    agenda_box.line.color.rgb = RGBColor(*COLORS["gray_border"])
    
    # Set margins
    agenda_tf.margin_left = Inches(0.3)
    agenda_tf.margin_top = Inches(0.2)
    agenda_tf.margin_right = Inches(0.3)
    agenda_tf.margin_bottom = Inches(0.2)
    
    # Add title
    agenda_title_p = agenda_tf.paragraphs[0]
    agenda_title_p.text = "Overview"
    agenda_title_p.font.bold = True
    agenda_title_p.font.size = Pt(FONT_SIZES["agenda_title"])
    agenda_title_p.font.color.rgb = RGBColor(*COLORS["black"])
    agenda_title_p.font.underline = True
    agenda_title_p.space_after = Pt(12)
    
    # Add agenda points
    for point in agenda_points:
        p = agenda_tf.add_paragraph()
        p.text = "❖ " + point
        p.font.size = Pt(FONT_SIZES["agenda_points"])
        p.font.color.rgb = RGBColor(*COLORS["black"])
        p.space_before = Pt(15)
        p.space_after = Pt(6)
