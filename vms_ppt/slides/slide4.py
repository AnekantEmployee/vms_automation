import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from .slide_utils import SlideUtils


def create_slide4(prs: Presentation, slide4_data):
    """Create slide(s) for Software Update and OS Hardening vulnerabilities with pagination"""
    start_time = time.time()
    
    # Check if we need to split into multiple slides
    software_vulns = slide4_data.get("table1", {}).get("rows", [])
    config_vulns = slide4_data.get("table2", {}).get("rows", [])
    
    slides_created = []
    
    # Define max rows per slide (accounting for table headers and layout)
    MAX_ROWS_PER_SLIDE = 12
    
    # Create Software Update/Uninstallation slides
    software_slides = _create_software_vulnerability_slides(prs, slide4_data, software_vulns, MAX_ROWS_PER_SLIDE)
    slides_created.extend(software_slides)
    
    # Create OS Hardening/Configuration slides  
    config_slides = _create_config_vulnerability_slides(prs, slide4_data, config_vulns, MAX_ROWS_PER_SLIDE)
    slides_created.extend(config_slides)
    
    # Create combined summary slide if we have multiple slides
    if len(slides_created) > 1:
        summary_slide = _create_category_summary_slide(prs, slide4_data, software_vulns, config_vulns)
        slides_created.insert(0, summary_slide)
    
    # Print results
    _print_slide4_results(start_time, software_vulns, config_vulns, len(slides_created))
    
    return slides_created, time.time() - start_time


def _create_software_vulnerability_slides(prs, slide4_data, software_vulns, max_rows):
    """Create one or more slides for software/application vulnerabilities"""
    if not software_vulns or len(software_vulns) == 0:
        return []
    
    slides_created = []
    software_chunks = list(_split_rows_for_slides(software_vulns, max_rows))
    
    for chunk_idx, chunk in enumerate(software_chunks):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Create title with page number if multiple slides
        title = f"{slide4_data['title']} - Software/Application Issues"
        if len(software_chunks) > 1:
            title += f" (Page {chunk_idx + 1} of {len(software_chunks)})"
        
        SlideUtils.create_title_bar(slide, prs, title)
        
        # Get layout parameters
        layout = SlideUtils.get_standard_layout_params(prs)
        
        # Create subtitle
        subtitle = f"Software Update/Uninstallation Vulnerabilities ({len(software_vulns)} total)"
        if len(software_chunks) > 1:
            start_idx = chunk_idx * max_rows + 1
            end_idx = min((chunk_idx + 1) * max_rows, len(software_vulns))
            subtitle += f" - Items {start_idx}-{end_idx}"
        
        SlideUtils.create_subtitle(slide, layout['table_left'], Inches(0.8), layout['max_table_width'], subtitle)
        
        # Calculate totals for this chunk
        chunk_totals = SlideUtils.calculate_column_totals(chunk, 6, "Page Total")
        
        # Create table
        data_with_totals = chunk + [chunk_totals]
        table = SlideUtils.create_table_with_headers(slide, len(data_with_totals) + 1, 6,
                                                   layout['table_left'], Inches(1.2), 
                                                   layout['max_table_width'], Inches(4.8))
        
        column_widths = [Inches(3.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(2.8)]
        SlideUtils.set_table_column_widths(table, column_widths)
        SlideUtils.set_table_row_heights(table, Inches(0.4), Inches(0.3))
        SlideUtils.format_header_row(table, slide4_data["table1"]["columns"])
        SlideUtils.populate_table_data(table, data_with_totals)
        
        # Add footer with navigation info
        if len(software_chunks) > 1:
            footer_text = f"Software Vulnerabilities - Page {chunk_idx + 1} of {len(software_chunks)}"
            SlideUtils.create_footnote(slide, layout['table_left'], Inches(6.2), 
                                     layout['max_table_width'], footer_text, font_size=10)
        
        slides_created.append(slide)
    
    return slides_created


def _create_config_vulnerability_slides(prs, slide4_data, config_vulns, max_rows):
    """Create one or more slides for configuration/hardening vulnerabilities"""
    if not config_vulns or len(config_vulns) == 0:
        return []
    
    slides_created = []
    config_chunks = list(_split_rows_for_slides(config_vulns, max_rows))
    
    for chunk_idx, chunk in enumerate(config_chunks):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Create title with page number if multiple slides
        title = f"{slide4_data['title']} - Configuration Issues"
        if len(config_chunks) > 1:
            title += f" (Page {chunk_idx + 1} of {len(config_chunks)})"
        
        SlideUtils.create_title_bar(slide, prs, title)
        
        # Get layout parameters
        layout = SlideUtils.get_standard_layout_params(prs)
        
        # Create subtitle
        subtitle = f"OS Hardening/Configuration Issues ({len(config_vulns)} total)"
        if len(config_chunks) > 1:
            start_idx = chunk_idx * max_rows + 1
            end_idx = min((chunk_idx + 1) * max_rows, len(config_vulns))
            subtitle += f" - Items {start_idx}-{end_idx}"
        
        SlideUtils.create_subtitle(slide, layout['table_left'], Inches(0.8), layout['max_table_width'], subtitle)
        
        # Calculate totals for this chunk
        chunk_totals = SlideUtils.calculate_column_totals(chunk, 4, "Page Total")
        
        # Create table
        data_with_totals = chunk + [chunk_totals]
        table = SlideUtils.create_table_with_headers(slide, len(data_with_totals) + 1, 4,
                                                   layout['table_left'], Inches(1.2), 
                                                   layout['max_table_width'], Inches(4.8))
        
        column_widths = [Inches(8.5), Inches(1.5), Inches(1.5), Inches(1.8)]
        SlideUtils.set_table_column_widths(table, column_widths)
        SlideUtils.set_table_row_heights(table, Inches(0.4), Inches(0.3))
        SlideUtils.format_header_row(table, slide4_data["table2"]["columns"])
        SlideUtils.populate_table_data(table, data_with_totals)
        
        # Add footer with navigation info
        if len(config_chunks) > 1:
            footer_text = f"Configuration Issues - Page {chunk_idx + 1} of {len(config_chunks)}"
            SlideUtils.create_footnote(slide, layout['table_left'], Inches(6.2), 
                                     layout['max_table_width'], footer_text, font_size=10)
        
        slides_created.append(slide)
    
    return slides_created


def _create_category_summary_slide(prs, slide4_data, software_vulns, config_vulns):
    """Create a summary overview slide when multiple detail slides exist"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Create title
    SlideUtils.create_title_bar(slide, prs, f"{slide4_data['title']} - Category Summary")
    
    # Get layout parameters
    layout = SlideUtils.get_standard_layout_params(prs)
    
    # Create summary statistics
    software_total = SlideUtils.calculate_column_totals(software_vulns, 6, "Total")
    config_total = SlideUtils.calculate_column_totals(config_vulns, 4, "Total")
    
    # Summary table
    summary_data = [
        ["Category", "Count", "Critical", "High", "Medium", "Low", "Total"],
        ["Software/Applications", str(len(software_vulns)), 
         str(software_total[1]), str(software_total[2]), str(software_total[3]), str(software_total[4]), str(software_total[5])],
        ["Configuration/Hardening", str(len(config_vulns)), 
         str(config_total[1]), "0", str(config_total[2]), "0", str(config_total[3])],
        ["Combined Total", str(len(software_vulns) + len(config_vulns)), 
         str(software_total[1] + config_total[1]), str(software_total[2]), 
         str(software_total[3] + config_total[2]), str(software_total[4]), 
         str(software_total[5] + config_total[3])]
    ]
    
    # Create summary table
    table = SlideUtils.create_table_with_headers(slide, 4, 7, layout['table_left'], Inches(1.2), 
                                               layout['max_table_width'], Inches(1.6))
    
    column_widths = [Inches(3.0), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.3)]
    SlideUtils.set_table_column_widths(table, column_widths)
    SlideUtils.set_table_row_heights(table, Inches(0.35), Inches(0.4))
    SlideUtils.format_header_row(table, summary_data[0])
    SlideUtils.populate_table_data(table, summary_data[1:])
    
    # Add key findings text box
    findings_text = f"""Key Findings:
• {len(software_vulns)} Software/Application vulnerabilities identified
• {len(config_vulns)} Configuration/Hardening issues found
• Detailed breakdown available in following slides
• Priority remediation required for critical and high severity items"""
    
    findings_box = slide.shapes.add_textbox(layout['table_left'], Inches(3.2), 
                                          layout['max_table_width'], Inches(2.0))
    findings_tf = findings_box.text_frame
    findings_tf.text = findings_text
    for paragraph in findings_tf.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
    
    return slide


def _split_rows_for_slides(rows, max_rows_per_slide):
    """Split list of rows into smaller chunks for pagination"""
    for i in range(0, len(rows), max_rows_per_slide):
        yield rows[i:i + max_rows_per_slide]


def _print_slide4_results(start_time, software_vulns, config_vulns, slides_created):
    """Print enhanced slide 4 results with pagination info"""
    runtime = time.time() - start_time
    print(f"Slide 4 series created successfully!")
    print(f"Software vulnerabilities: {len(software_vulns)} items")
    print(f"Configuration vulnerabilities: {len(config_vulns)} items") 
    print(f"Total slides created: {slides_created}")
    
    if len(software_vulns) > 12:
        print(f"✓ Software vulnerabilities split across multiple slides for readability")
    if len(config_vulns) > 12:
        print(f"✓ Configuration vulnerabilities split across multiple slides for readability")
        
    print(f"Runtime: {runtime:.4f} seconds")


# Handle slide 4 creation in main function
def handle_slide4_creation(prs, slide4_data):
    """Handle slide 4 creation which might return multiple slides"""
    result = create_slide4(prs, slide4_data)
    
    if isinstance(result, tuple):
        slides_created, runtime = result
        return slides_created, runtime
    else:
        # Fallback for single slide
        return [result], 0
