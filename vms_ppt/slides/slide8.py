import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from .slide_utils import SlideUtils


def create_slide8(prs: Presentation, slide8_data):
    """Create slide(s) for High Priority Vulnerability Response with pagination"""
    start_time = time.time()
    
    # Check if we need to split into multiple slides
    high_data = slide8_data.get("table", {}).get("rows", [])
    
    slides_created = []
    
    # Define max rows per slide (accounting for table headers and layout)
    MAX_ROWS_PER_SLIDE = 12  # Same as slide 7
    
    # Filter out existing total rows and create pagination
    filtered_rows = [row for row in high_data if "Total" not in str(row[0])]
    
    # Create paginated slides for high priority data
    high_slides = _create_high_slides(prs, slide8_data, filtered_rows, MAX_ROWS_PER_SLIDE)
    slides_created.extend(high_slides)
    
    # Create summary slide if we have multiple slides
    if len(slides_created) > 1:
        summary_slide = _create_high_summary_slide(prs, slide8_data, filtered_rows)
        slides_created.insert(0, summary_slide)
    
    # Print results
    _print_slide8_results(start_time, filtered_rows, len(slides_created))
    
    return slides_created, time.time() - start_time


def _create_high_slides(prs, slide8_data, high_data, max_rows):
    """Create one or more slides for high priority vulnerability data with optimized column widths"""
    if not high_data or len(high_data) == 0:
        return []
    
    slides_created = []
    high_chunks = list(_split_rows_for_slides(high_data, max_rows))
    
    for chunk_idx, chunk in enumerate(high_chunks):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Create title with page number if multiple slides
        title = slide8_data["title"]
        if len(high_chunks) > 1:
            title += f" (Page {chunk_idx + 1} of {len(high_chunks)})"
        
        SlideUtils.create_title_bar(slide, prs, title)
        
        # Get layout parameters
        layout = SlideUtils.get_standard_layout_params(prs)
        
        # Create subtitle with item range
        subtitle = f"High Priority Response Items ({len(high_data)} total items)"
        if len(high_chunks) > 1:
            start_idx = chunk_idx * max_rows + 1
            end_idx = min((chunk_idx + 1) * max_rows, len(high_data))
            subtitle += f" - Items {start_idx}-{end_idx}"
        
        SlideUtils.create_subtitle(slide, layout['table_left'], Inches(0.8), 
                                 layout['max_table_width'], subtitle)
        
        # Calculate totals for this chunk (8 columns)
        chunk_totals = SlideUtils.calculate_column_totals(chunk, 8, "Page Total")
        
        # Create table (8 columns matching slide7 structure)
        data_with_totals = chunk + [chunk_totals]
        table = SlideUtils.create_table_with_headers(slide, len(data_with_totals) + 1, 
                                                   8,  # 8 columns total
                                                   layout['table_left'], Inches(1.2), 
                                                   layout['max_table_width'], Inches(4.8))
        
        # Set column widths matching slide7 structure (8 columns)
        # Asset & Network Location: 20%, Vulnerability Description: 20%
        # Critical/High/Medium/Low: 4% each (16% total)
        # CVE IDs: 20%, Business Impact: 20%
        column_widths = [
            Inches(2.5),   # Asset & Network Location (20%)
            Inches(2.5),   # Vulnerability Description (20%)
            Inches(0.5),   # Critical (4%)
            Inches(0.5),   # High (4%)
            Inches(0.5),   # Medium (4%)
            Inches(0.5),   # Low (4%)
            Inches(2.5),   # CVE IDs & CVSS Score (20%)
            Inches(3.0)    # Business Impact Context (24%)
        ]
        
        # Use safe column width setting to avoid TypeError
        _set_column_widths_safe(table, column_widths)
        
        SlideUtils.set_table_row_heights(table, Inches(0.35), Inches(0.32))
        SlideUtils.format_header_row(table, slide8_data["table"]["columns"])
        SlideUtils.populate_table_data(table, data_with_totals)
        
        # Add footer with navigation info
        if len(high_chunks) > 1:
            footer_text = f"High Priority Response - Page {chunk_idx + 1} of {len(high_chunks)}"
            SlideUtils.create_footnote(slide, layout['table_left'], Inches(6.2), 
                                     layout['max_table_width'], footer_text, font_size=10)
        
        slides_created.append(slide)
    
    return slides_created


def _create_high_summary_slide(prs, slide8_data, high_data):
    """Create a summary overview slide when multiple detail slides exist"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Create title
    SlideUtils.create_title_bar(slide, prs, f"{slide8_data['title']} - Executive Summary")
    
    # Get layout parameters
    layout = SlideUtils.get_standard_layout_params(prs)
    
    # Calculate overall totals (adjusted for 8 columns)
    overall_totals = SlideUtils.calculate_column_totals(high_data, 8, "Grand Total")
    
    # Create summary statistics table
    summary_data = [
        ["Metric", "Value"],
        ["Total High Priority Assets", str(len(high_data))],
        ["Critical Vulnerabilities", str(overall_totals[2]) if len(overall_totals) > 2 else "0"],
        ["High Priority Vulnerabilities", str(overall_totals[3]) if len(overall_totals) > 3 else "0"],
        ["Medium Priority Vulnerabilities", str(overall_totals[4]) if len(overall_totals) > 4 else "0"],
        ["Low Priority Vulnerabilities", str(overall_totals[5]) if len(overall_totals) > 5 else "0"],
        ["CVE IDs Identified", "Multiple - See Details"],
        ["Business Impact Cases", str(len([row for row in high_data if row[7] != "Standard Risk"]))]
    ]
    
    # Create summary table
    table = SlideUtils.create_table_with_headers(slide, len(summary_data), 2, 
                                               layout['table_left'], Inches(1.2), 
                                               Inches(8.0), Inches(2.8))
    
    column_widths = [Inches(4.0), Inches(4.0)]
    SlideUtils.set_table_column_widths(table, column_widths)
    SlideUtils.set_table_row_heights(table, Inches(0.35), Inches(0.35))
    SlideUtils.format_header_row(table, summary_data[0])
    SlideUtils.populate_table_data(table, summary_data[1:])
    
    # Add key findings text box
    findings_text = f"""Key Findings - High Priority Response:
• {len(high_data)} high priority assets requiring prompt remediation
• Multiple CVE IDs identified with varying CVSS scores
• Comprehensive business impact assessment completed
• Asset-specific network location and service context provided
• Detailed breakdown available in following slides for operational teams"""
    
    findings_box = slide.shapes.add_textbox(layout['table_left'], Inches(4.2), 
                                          layout['max_table_width'], Inches(1.8))
    findings_tf = findings_box.text_frame
    findings_tf.text = findings_text
    for paragraph in findings_tf.paragraphs:
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
    
    return slide


def _set_column_widths_safe(table, widths):
    """Safely set column widths ensuring they are integers"""
    for i, width in enumerate(widths):
        if i < len(table.columns):
            # Convert to int to avoid TypeError
            if hasattr(width, 'value'):  # Handle Inches objects
                width_value = int(width)
            else:
                width_value = int(width)
            table.columns[i].width = width_value


def _split_rows_for_slides(rows, max_rows_per_slide):
    """Split list of rows into smaller chunks for pagination"""
    for i in range(0, len(rows), max_rows_per_slide):
        yield rows[i:i + max_rows_per_slide]


def _print_slide8_results(start_time, high_data, slides_created):
    """Print enhanced slide 8 results with pagination info"""
    runtime = time.time() - start_time
    print(f"Slide 8 series created successfully!")
    print(f"High priority vulnerability items: {len(high_data)} assets")
    print(f"Total slides created: {slides_created}")
    
    if len(high_data) > 12:
        print(f"✓ High priority vulnerability data split across multiple slides for readability")
        
    print(f"Runtime: {runtime:.4f} seconds")
