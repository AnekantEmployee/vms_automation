from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
import time
from config_colors import COLORS, FONT_SIZES

def create_slide2(prs: Presentation, slide2_data):
    start_time = time.time()
    
    # Calculate vulnerability summary totals dynamically
    vuln_totals = ["Grand Total", 0, 0, 0, 0, 0, 0]
    
    for row in slide2_data["vulnerability_summary"]["rows"]:
        for col_idx in range(1, 7):  # Skip category name column
            value = row[col_idx].strip() if row[col_idx] else "0"
            if value.isdigit():
                vuln_totals[col_idx] += int(value)
    
    # Calculate grand total
    vuln_totals[6] = sum(vuln_totals[1:6])
    
    # Calculate baselining failures total dynamically
    failures_total = 0
    for row in slide2_data["baselining_failures"]["rows"]:
        if "Total" not in row[0]:  # Skip any existing total rows
            value = row[1].strip() if row[1] else "0"
            if value.isdigit():
                failures_total += int(value)
    
    failures_total_row = ["Total", str(failures_total)]
    
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Create blue title bar
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = Inches(0.6)

    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS["blue"])
    shape.line.color.rgb = RGBColor(*COLORS["blue"])

    # Add title text box
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = slide2_data["title"]
    p.font.bold = True
    p.font.size = Pt(FONT_SIZES["title"])
    p.font.color.rgb = RGBColor(*COLORS["white"])
    p.alignment = 1

    # Layout parameters
    table_left = Inches(0.2)
    max_table_width = prs.slide_width - Inches(0.4)
    
    # Left column tables (Baseline Status and Baselining Failures)
    left_table_width = Inches(4.8)
    
    # Table 1: Baseline Status
    baseline_top = Inches(0.8)
    
    # Baseline status title
    baseline_title_box = slide.shapes.add_textbox(table_left, baseline_top, left_table_width, Inches(0.25))
    baseline_tf = baseline_title_box.text_frame
    baseline_p = baseline_tf.paragraphs[0]
    baseline_p.text = slide2_data["baseline_status"]["title"]
    baseline_p.font.bold = True
    baseline_p.font.size = Pt(12)
    baseline_p.font.color.rgb = RGBColor(*COLORS["white"])
    baseline_p.alignment = 1
    baseline_title_box.fill.solid()
    baseline_title_box.fill.fore_color.rgb = RGBColor(*COLORS["blue"])

    # Baseline status table
    baseline_table = slide.shapes.add_table(2, 2, table_left, Inches(1.05), left_table_width, Inches(1.0)).table
    baseline_table.columns[0].width = Inches(3.5)
    baseline_table.columns[1].width = Inches(1.3)

    # Headers
    for i, col_name in enumerate(slide2_data["baseline_status"]["columns"]):
        cell = baseline_table.cell(0, i)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*COLORS["blue"])
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(*COLORS["white"])
        p.alignment = 1

    # Data
    for col_idx, cell_data in enumerate(slide2_data["baseline_status"]["rows"][0]):
        cell = baseline_table.cell(1, col_idx)
        cell.text = cell_data
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        if col_idx > 0:
            p.alignment = 1

    # Footnote
    footnote1_box = slide.shapes.add_textbox(table_left, Inches(2.1), left_table_width, Inches(0.2))
    footnote1_tf = footnote1_box.text_frame
    footnote1_p = footnote1_tf.paragraphs[0]
    footnote1_p.text = slide2_data["baseline_status"]["footnote"]
    footnote1_p.font.size = Pt(8)
    footnote1_p.font.italic = True

    # Table 2: Baselining Failures with dynamic total
    failures_top = Inches(2.4)
    
    # Failures title
    failures_title_box = slide.shapes.add_textbox(table_left, failures_top, left_table_width, Inches(0.25))
    failures_tf = failures_title_box.text_frame
    failures_p = failures_tf.paragraphs[0]
    failures_p.text = slide2_data["baselining_failures"]["title"]
    failures_p.font.bold = True
    failures_p.font.size = Pt(12)
    failures_p.font.color.rgb = RGBColor(*COLORS["white"])
    failures_p.alignment = 1
    failures_title_box.fill.solid()
    failures_title_box.fill.fore_color.rgb = RGBColor(*COLORS["blue"])

    # Failures table with calculated total
    failures_data_with_total = slide2_data["baselining_failures"]["rows"] + [failures_total_row]
    failures_rows = len(failures_data_with_total)
    failures_table = slide.shapes.add_table(failures_rows, 2, table_left, Inches(2.65), left_table_width, Inches(1.1)).table
    failures_table.columns[0].width = Inches(3.5)
    failures_table.columns[1].width = Inches(1.3)

    for row_idx, row_data in enumerate(failures_data_with_total):
        for col_idx, cell_data in enumerate(row_data):
            cell = failures_table.cell(row_idx, col_idx)
            cell.text = cell_data
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            
            # Blue background for Total row
            if cell_data == "Total" or (row_data[0] == "Total" and col_idx == 1):
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*COLORS["blue"])
                p.font.bold = True
                p.font.color.rgb = RGBColor(*COLORS["white"])
            
            if col_idx > 0:
                p.alignment = 1

    # Footnote 2
    footnote2_box = slide.shapes.add_textbox(table_left, Inches(3.8), left_table_width, Inches(0.3))
    footnote2_tf = footnote2_box.text_frame
    footnote2_p = footnote2_tf.paragraphs[0]
    footnote2_p.text = slide2_data["baselining_failures"]["footnote"]
    footnote2_p.font.size = Pt(8)
    footnote2_p.font.italic = True

    # Right side: Vulnerability Summary Table with dynamic totals
    right_table_left = Inches(5.2)
    right_table_width = prs.slide_width - right_table_left - Inches(0.2)
    
    vuln_data_with_totals = slide2_data["vulnerability_summary"]["rows"] + [vuln_totals]
    vuln_rows = len(vuln_data_with_totals) + 1  # +1 for header
    vuln_cols = len(slide2_data["vulnerability_summary"]["columns"])
    
    vuln_table = slide.shapes.add_table(vuln_rows, vuln_cols, right_table_left, Inches(0.8), right_table_width, Inches(1.5)).table
    
    # Set column widths to fit properly
    vuln_table.columns[0].width = Inches(2.5)
    for i in range(1, 6):
        vuln_table.columns[i].width = Inches(0.9)
    vuln_table.columns[6].width = Inches(1.0)

    # Header row
    for i, col_name in enumerate(slide2_data["vulnerability_summary"]["columns"]):
        cell = vuln_table.cell(0, i)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*COLORS["blue"])
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(*COLORS["white"])
        p.alignment = 1

    # Data rows
    for row_idx, row_data in enumerate(vuln_data_with_totals):
        for col_idx, cell_data in enumerate(row_data):
            cell = vuln_table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data) if cell_data else ""
            
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(*COLORS["black"])
            
            if col_idx > 0:
                p.alignment = 1
            
            # Blue background for Grand Total row
            if row_data[0] == "Grand Total":
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*COLORS["blue"])
                p.font.bold = True
                p.font.color.rgb = RGBColor(*COLORS["white"])
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*COLORS["very_light_gray"])

    # Legends section
    legends_top = Inches(2.5)
    legends_box = slide.shapes.add_textbox(right_table_left, legends_top, right_table_width, Inches(1.0))
    legends_tf = legends_box.text_frame
    
    # Legends title
    legends_title_p = legends_tf.paragraphs[0]
    legends_title_p.text = "Legends"
    legends_title_p.font.bold = True
    legends_title_p.font.size = Pt(12)
    
    for legend in slide2_data["legends"]:
        legend_p = legends_tf.add_paragraph()
        legend_p.text = f"• {legend['term']}: {legend['description']}"
        legend_p.font.size = Pt(8)
        legend_p.space_before = Pt(3)

    # Timeline table
    timeline_top = Inches(3.8)
    timeline_rows = len(slide2_data["timeline"]["rows"]) + 1
    timeline_cols = len(slide2_data["timeline"]["columns"])
    
    timeline_table = slide.shapes.add_table(timeline_rows, timeline_cols, table_left, timeline_top, max_table_width, Inches(2.8)).table
    
    # Set timeline column widths
    timeline_table.columns[0].width = Inches(1.6)   # Severity
    timeline_table.columns[1].width = Inches(2.0)   # CVSS Score
    timeline_table.columns[2].width = Inches(2.2)   # Timeline
    timeline_table.columns[3].width = Inches(6.5)   # Description

    # Timeline header
    for i, col_name in enumerate(slide2_data["timeline"]["columns"]):
        cell = timeline_table.cell(0, i)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*COLORS["blue"])
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(*COLORS["white"])
        p.alignment = 1

    # Timeline data with color coding
    severity_colors = {
        "Immediate": (255, 0, 0),     # Red
        "Critical": (220, 20, 60),    # Crimson  
        "High": (255, 165, 0),        # Orange
        "Medium": (255, 255, 0),      # Yellow
        "Low": (128, 128, 128)        # Gray
    }

    for row_idx, row_data in enumerate(slide2_data["timeline"]["rows"]):
        severity = row_data[0]
        for col_idx, cell_data in enumerate(row_data):
            cell = timeline_table.cell(row_idx + 1, col_idx)
            cell.text = cell_data
            
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            
            if col_idx == 0:  # Severity column
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*severity_colors[severity])
                p.font.bold = True
                p.font.color.rgb = RGBColor(*COLORS["white"])
                p.alignment = 1

    # Disclaimer
    disclaimer_box = slide.shapes.add_textbox(table_left, Inches(6.8), max_table_width, Inches(0.3))
    disclaimer_tf = disclaimer_box.text_frame
    disclaimer_p = disclaimer_tf.paragraphs[0]
    disclaimer_p.text = f"Disclaimer – {slide2_data['disclaimer']}"
    disclaimer_p.font.size = Pt(8)
    disclaimer_p.font.italic = True

    end_time = time.time()
    runtime = end_time - start_time
    
    print(f"Slide 2 created successfully!")
    print(f"Baselining Failures Total: {failures_total}")
    print(f"Vulnerability Summary Totals - Immediate: {vuln_totals[1]}, Critical: {vuln_totals[2]}, High: {vuln_totals[3]}, Medium: {vuln_totals[4]}, Low: {vuln_totals[5]}, Grand Total: {vuln_totals[6]}")
    print(f"Runtime: {runtime:.4f} seconds")
    
    return runtime
