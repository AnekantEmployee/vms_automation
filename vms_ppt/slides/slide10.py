import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from .slide_utils import SlideUtils


def create_slide10(prs: Presentation, slide10_data):
    """Create executive summary and strategic roadmap slide with black text"""
    start_time = time.time()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Create title bar
    SlideUtils.create_title_bar(slide, prs, slide10_data["title"])
    
    # Get layout parameters
    layout = SlideUtils.get_standard_layout_params(prs)
    
    # Create subtitle
    SlideUtils.create_subtitle(slide, layout['table_left'], Inches(0.8), 
                             layout['max_table_width'], slide10_data.get("subtitle", ""))
    
    # Create main sections
    current_y = Inches(1.3)
    
    for section in slide10_data.get("sections", []):
        # Section title - CHANGED FROM BLUE TO BLACK
        section_title = slide.shapes.add_textbox(layout['table_left'], current_y, 
                                               layout['max_table_width'], Inches(0.3))
        section_tf = section_title.text_frame
        section_p = section_tf.paragraphs[0]
        section_p.text = section["title"]
        section_p.font.bold = True
        section_p.font.size = Pt(14)
        section_p.font.color.rgb = RGBColor(0, 0, 0)  # Changed from (0, 0, 139) to black
        current_y += Inches(0.4)
        
        # Create table for section
        table = SlideUtils.create_table_with_headers(slide, len(section["rows"]) + 1, 
                                                   len(section["columns"]),
                                                   layout['table_left'], current_y, 
                                                   layout['max_table_width'], Inches(1.8))
        
        # Set column widths based on section type
        if section["type"] == "metrics":
            column_widths = [Inches(4.0), Inches(2.5), Inches(2.5)]
        elif section["type"] == "action_plan":
            column_widths = [Inches(3.0), Inches(2.5), Inches(2.0), Inches(3.5)]
        else:  # insights
            column_widths = [Inches(3.0), Inches(2.5), Inches(5.5)]
        
        SlideUtils.set_table_column_widths(table, column_widths)
        SlideUtils.set_table_row_heights(table, Inches(0.3), Inches(0.25))
        SlideUtils.format_header_row(table, section["columns"])
        SlideUtils.populate_table_data(table, section["rows"])
        
        current_y += Inches(2.2)
    
    # Add key recommendations box - CHANGED FROM BLUE TO BLACK
    recommendations_y = Inches(5.8)
    rec_title = slide.shapes.add_textbox(layout['table_left'], recommendations_y, 
                                       Inches(6.0), Inches(0.3))
    rec_tf = rec_title.text_frame
    rec_p = rec_tf.paragraphs[0]
    rec_p.text = "Strategic Recommendations"
    rec_p.font.bold = True
    rec_p.font.size = Pt(14)
    rec_p.font.color.rgb = RGBColor(0, 0, 0)  # Changed from (0, 0, 139) to black
    
    rec_box = slide.shapes.add_textbox(layout['table_left'], recommendations_y + Inches(0.35), 
                                     Inches(6.0), Inches(1.0))
    rec_text_frame = rec_box.text_frame
    rec_text_frame.text = "• " + "\n• ".join(slide10_data.get("key_recommendations", [])[:4])
    for paragraph in rec_text_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Already black, kept as is
    
    # Add success metrics box - CHANGED FROM BLUE TO BLACK
    metrics_box = slide.shapes.add_textbox(Inches(6.5), recommendations_y, 
                                         Inches(6.0), Inches(1.35))
    metrics_tf = metrics_box.text_frame
    metrics_p = metrics_tf.paragraphs[0]
    metrics_p.text = "Success Metrics"
    metrics_p.font.bold = True
    metrics_p.font.size = Pt(14)
    metrics_p.font.color.rgb = RGBColor(0, 0, 0)  # Changed from (0, 0, 139) to black
    
    for metric in slide10_data.get("success_metrics", [])[:4]:
        metric_p = metrics_tf.add_paragraph()
        metric_p.text = f"• {metric}"
        metric_p.font.size = Pt(10)
        metric_p.font.color.rgb = RGBColor(0, 0, 0)  # Already black, kept as is
    
    # Print results
    runtime = time.time() - start_time
    print(f"Slide 10 created successfully!")
    print(f"Executive summary with {len(slide10_data.get('sections', []))} strategic sections")
    print(f"Runtime: {runtime:.4f} seconds")
    
    return runtime
