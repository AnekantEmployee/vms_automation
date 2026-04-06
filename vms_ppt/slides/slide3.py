import time
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx import Presentation
from .slide_utils import SlideUtils


def create_slide3(prs: Presentation, slide3_data):
    """Create slide(s) for Critical & High Risk Vulnerabilities with pagination (max 15 rows per slide)"""
    start_time = time.time()
    
    # Check if we need to split into multiple slides
    critical_vulns = slide3_data.get("table1", {}).get("rows", [])
    high_vulns = slide3_data.get("table2", {}).get("rows", [])
    
    slides_created = []
    
    # Define max rows per slide - SET TO 15 AS REQUESTED
    MAX_ROWS_PER_SLIDE = 15
    
    # Create Critical Vulnerability slides with pagination
    critical_slides = _create_critical_vulnerability_slides(prs, slide3_data, critical_vulns, MAX_ROWS_PER_SLIDE)
    slides_created.extend(critical_slides)
    
    # Create High Vulnerability slides with pagination
    high_slides = _create_high_vulnerability_slides(prs, slide3_data, high_vulns, MAX_ROWS_PER_SLIDE)
    slides_created.extend(high_slides)
    
    # Create combined summary slide if we have multiple slides
    if len(slides_created) > 1:
        summary_slide = _create_vulnerability_summary_slide(prs, slide3_data, critical_vulns, high_vulns)
        slides_created.insert(0, summary_slide)  # Add summary as first slide
    
    # Print results
    _print_slide3_results(start_time, critical_vulns, high_vulns, len(slides_created))
    
    return slides_created, time.time() - start_time


def _create_critical_vulnerability_slides(prs, slide3_data, critical_vulns, max_rows):
    """Create one or more slides for critical vulnerabilities (max 15 per slide)"""
    if not critical_vulns or len(critical_vulns) == 0:
        return []
    
    slides_created = []
    critical_chunks = list(_split_rows_for_slides(critical_vulns, max_rows))
    
    for chunk_idx, chunk in enumerate(critical_chunks):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Create title - add page number if multiple slides
        title = f"{slide3_data['title']} - Critical Vulnerabilities"
        if len(critical_chunks) > 1:
            title += f" (Page {chunk_idx + 1} of {len(critical_chunks)})"
        
        SlideUtils.create_title_bar(slide, prs, title)
        
        # Get layout parameters
        layout = SlideUtils.get_standard_layout_params(prs)
        
        # Create subtitle with item range
        subtitle = f"Critical Severity Vulnerabilities ({len(critical_vulns)} total)"
        if len(critical_chunks) > 1:
            start_idx = chunk_idx * max_rows + 1
            end_idx = min((chunk_idx + 1) * max_rows, len(critical_vulns))
            subtitle += f" - Items {start_idx}-{end_idx}"
        
        SlideUtils.create_subtitle(slide, layout['table_left'], Inches(0.8), layout['max_table_width'], subtitle)
        
        # Calculate totals for this chunk (only for Total CVEs column - index 2)
        total_cves = sum(int(row[2]) if row[2].isdigit() else 0 for row in chunk)
        chunk_totals = ["Page Total", "", str(total_cves)]
        
        # Create table with 3 columns
        data_with_totals = chunk + [chunk_totals]
        table = SlideUtils.create_table_with_headers(slide, len(data_with_totals) + 1, 3,
                                                   layout['table_left'], Inches(1.2), 
                                                   layout['max_table_width'], Inches(4.8))
        
        # Updated column widths for 3 columns: Title (wider), CVE IDs (wider), Total CVEs (smaller)
        column_widths = [Inches(6.0), Inches(4.8), Inches(1.5)]
        SlideUtils.set_table_column_widths(table, column_widths)
        SlideUtils.set_table_row_heights(table, Inches(0.35), Inches(0.32))
        SlideUtils.format_header_row(table, slide3_data["table1"]["columns"])
        SlideUtils.populate_table_data(table, data_with_totals)
        
        # Add footer with navigation info for multiple pages
        if len(critical_chunks) > 1:
            footer_text = f"Critical Vulnerabilities - Page {chunk_idx + 1} of {len(critical_chunks)} (Max {max_rows} items per page)"
            SlideUtils.create_footnote(slide, layout['table_left'], Inches(6.2), 
                                     layout['max_table_width'], footer_text, font_size=10)
        
        slides_created.append(slide)
    
    return slides_created


def _create_high_vulnerability_slides(prs, slide3_data, high_vulns, max_rows):
    """Create one or more slides for high vulnerabilities (max 15 per slide)"""
    if not high_vulns or len(high_vulns) == 0:
        return []
    
    slides_created = []
    high_chunks = list(_split_rows_for_slides(high_vulns, max_rows))
    
    for chunk_idx, chunk in enumerate(high_chunks):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Create title - add page number if multiple slides
        title = f"{slide3_data['title']} - High Vulnerabilities"
        if len(high_chunks) > 1:
            title += f" (Page {chunk_idx + 1} of {len(high_chunks)})"
        
        SlideUtils.create_title_bar(slide, prs, title)
        
        # Get layout parameters
        layout = SlideUtils.get_standard_layout_params(prs)
        
        # Create subtitle with item range
        subtitle = f"High Severity Vulnerabilities ({len(high_vulns)} total)"
        if len(high_chunks) > 1:
            start_idx = chunk_idx * max_rows + 1
            end_idx = min((chunk_idx + 1) * max_rows, len(high_vulns))
            subtitle += f" - Items {start_idx}-{end_idx}"
        
        SlideUtils.create_subtitle(slide, layout['table_left'], Inches(0.8), layout['max_table_width'], subtitle)
        
        # Calculate totals for this chunk (only for Total CVEs column - index 2)
        total_cves = sum(int(row[2]) if row[2].isdigit() else 0 for row in chunk)
        chunk_totals = ["Page Total", "", str(total_cves)]
        
        # Create table with 3 columns
        data_with_totals = chunk + [chunk_totals]
        table = SlideUtils.create_table_with_headers(slide, len(data_with_totals) + 1, 3,
                                                   layout['table_left'], Inches(1.2), 
                                                   layout['max_table_width'], Inches(4.8))
        
        # Updated column widths for 3 columns: Title (wider), CVE IDs (wider), Total CVEs (smaller)
        column_widths = [Inches(6.0), Inches(4.8), Inches(1.5)]
        SlideUtils.set_table_column_widths(table, column_widths)
        SlideUtils.set_table_row_heights(table, Inches(0.35), Inches(0.32))
        SlideUtils.format_header_row(table, slide3_data["table2"]["columns"])
        SlideUtils.populate_table_data(table, data_with_totals)
        
        # Add footer with navigation info for multiple pages
        if len(high_chunks) > 1:
            footer_text = f"High Vulnerabilities - Page {chunk_idx + 1} of {len(high_chunks)} (Max {max_rows} items per page)"
            SlideUtils.create_footnote(slide, layout['table_left'], Inches(6.2), 
                                     layout['max_table_width'], footer_text, font_size=10)
        
        slides_created.append(slide)
    
    return slides_created


def _create_vulnerability_summary_slide(prs, slide3_data, critical_vulns, high_vulns):
    """Create a summary overview slide when multiple detail slides exist"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Create title
    SlideUtils.create_title_bar(slide, prs, f"{slide3_data['title']} - Executive Summary")
    
    # Get layout parameters
    layout = SlideUtils.get_standard_layout_params(prs)
    
    # Calculate total CVEs for each category
    critical_total_cves = sum(int(row[2]) if row[2].isdigit() else 0 for row in critical_vulns)
    high_total_cves = sum(int(row[2]) if row[2].isdigit() else 0 for row in high_vulns)
    
    # Calculate number of pages needed
    max_rows = 15  # Updated to match MAX_ROWS_PER_SLIDE
    critical_pages = (len(critical_vulns) + max_rows - 1) // max_rows  # Ceiling division
    high_pages = (len(high_vulns) + max_rows - 1) // max_rows
    
    # Summary table
    summary_data = [
        ["Vulnerability Category", "Count", "Pages", "Total CVEs"],
        ["Critical Severity", str(len(critical_vulns)), str(critical_pages), str(critical_total_cves)],
        ["High Severity", str(len(high_vulns)), str(high_pages), str(high_total_cves)],
        ["Combined Total", str(len(critical_vulns) + len(high_vulns)), str(critical_pages + high_pages), str(critical_total_cves + high_total_cves)]
    ]
    
    # Create summary table with 4 columns
    table = SlideUtils.create_table_with_headers(slide, 4, 4, layout['table_left'], Inches(1.2), 
                                               layout['max_table_width'], Inches(1.6))
    
    column_widths = [Inches(4.0), Inches(2.0), Inches(2.0), Inches(2.3)]
    SlideUtils.set_table_column_widths(table, column_widths)
    SlideUtils.set_table_row_heights(table, Inches(0.35), Inches(0.4))
    SlideUtils.format_header_row(table, summary_data[0])
    SlideUtils.populate_table_data(table, summary_data[1:])
    
    # Add risk impact information
    if slide3_data.get("footnote"):
        _create_risk_impact_table(slide, slide3_data, layout, top_position=Inches(3.2))
    
    # Add key findings text box
    findings_text = f"""Key Findings ({max_rows} items max per slide):
• {len(critical_vulns)} Critical vulnerabilities requiring immediate attention ({critical_pages} pages)
• {len(high_vulns)} High severity vulnerabilities for priority remediation ({high_pages} pages)
• Total of {critical_pages + high_pages} detail slides for comprehensive review
• Total CVEs: {critical_total_cves + high_total_cves} across all vulnerabilities
• Each slide limited to {max_rows} vulnerabilities for optimal readability"""
    
    findings_box = slide.shapes.add_textbox(layout['table_left'], Inches(4.5), 
                                          layout['max_table_width'], Inches(1.5))
    findings_tf = findings_box.text_frame
    findings_tf.text = findings_text
    for paragraph in findings_tf.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
    
    return slide


def _create_risk_impact_table(slide, slide3_data, layout, top_position=Inches(6.4)):
    """Create risk and impact table with flexible positioning"""
    if not slide3_data.get("footnote"):
        return
        
    table_data = {
        "columns": ["Risk", "Impact"],
        "rows": [[slide3_data["footnote"]["Risk"], slide3_data["footnote"]["Impact"]]]
    }
    
    table = SlideUtils.create_table_with_headers(slide, 2, 2, layout['table_left'], top_position, 
                                               layout['max_table_width'], Inches(0.8))
    SlideUtils.set_table_column_widths(table, [Inches(6.165), Inches(6.165)])
    SlideUtils.set_table_row_heights(table, Inches(0.35), Inches(0.4))
    SlideUtils.format_header_row(table, table_data["columns"])
    SlideUtils.populate_table_data(table, table_data["rows"])


def _split_rows_for_slides(rows, max_rows_per_slide):
    """Split list of rows into smaller chunks for pagination"""
    for i in range(0, len(rows), max_rows_per_slide):
        yield rows[i:i + max_rows_per_slide]


def _print_slide3_results(start_time, critical_vulns, high_vulns, slides_created):
    """Print enhanced slide 3 results with pagination info"""
    runtime = time.time() - start_time
    print(f"Slide 3 series created successfully!")
    print(f"Critical vulnerabilities: {len(critical_vulns)} items")
    print(f"High severity vulnerabilities: {len(high_vulns)} items") 
    print(f"Total slides created: {slides_created}")
    
    # Show pagination info
    max_rows = 15  # Updated to match MAX_ROWS_PER_SLIDE
    if len(critical_vulns) > max_rows:
        critical_pages = (len(critical_vulns) + max_rows - 1) // max_rows
        print(f"✓ Critical vulnerabilities split across {critical_pages} slides (max {max_rows} per slide)")
    if len(high_vulns) > max_rows:
        high_pages = (len(high_vulns) + max_rows - 1) // max_rows
        print(f"✓ High vulnerabilities split across {high_pages} slides (max {max_rows} per slide)")
        
    print(f"Runtime: {runtime:.4f} seconds")
