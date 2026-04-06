from pptx import Presentation
import time
from pptx.util import Inches, Pt
from .slide_utils import SlideUtils
from pptx.dml.color import RGBColor


def create_slide7(prs: Presentation, slide7_data):
    """Create slide(s) for Stage 1 Critical Vulnerability Response with business impact column"""
    start_time = time.time()
    
    # Check if we need to split into multiple slides
    patching_data = slide7_data.get("table", {}).get("rows", [])
    
    slides_created = []
    
    # Define max rows per slide (accounting for table headers and layout)
    MAX_ROWS_PER_SLIDE = 10
    
    # Filter out existing total rows and create pagination
    filtered_rows = [row for row in patching_data if "Total" not in str(row[0])]
    
    # Create paginated slides for patching data
    patching_slides = _create_patching_slides(prs, slide7_data, filtered_rows, MAX_ROWS_PER_SLIDE)
    slides_created.extend(patching_slides)
    
    # Create summary slide if we have multiple slides
    if len(slides_created) > 1:
        summary_slide = _create_patching_summary_slide(prs, slide7_data, filtered_rows)
        slides_created.insert(0, summary_slide)
    
    # Print results
    _print_slide7_results(start_time, filtered_rows, len(slides_created))
    
    return slides_created, time.time() - start_time

def _create_patching_slides(prs, slide7_data, patching_data, max_rows):
    """Create one or more slides for critical vulnerability data with business impact column"""
    if not patching_data or len(patching_data) == 0:
        return []
    
    slides_created = []
    patching_chunks = list(_split_rows_for_slides(patching_data, max_rows))
    
    for chunk_idx, chunk in enumerate(patching_chunks):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Create title with page number if multiple slides
        title = slide7_data["title"]
        if len(patching_chunks) > 1:
            title += f" (Page {chunk_idx + 1} of {len(patching_chunks)})"
        
        SlideUtils.create_title_bar(slide, prs, title)
        
        # Get layout parameters
        layout = SlideUtils.get_standard_layout_params(prs)
        
        # Create subtitle with item range
        subtitle = f"Critical Vulnerability Response Items ({len(patching_data)} total items)"
        if len(patching_chunks) > 1:
            start_idx = chunk_idx * max_rows + 1
            end_idx = min((chunk_idx + 1) * max_rows, len(patching_data))
            subtitle += f" - Items {start_idx}-{end_idx}"
        
        SlideUtils.create_subtitle(slide, layout['table_left'], Inches(0.8), 
                                 layout['max_table_width'], subtitle)
        
        # Calculate totals for this chunk (8 columns total including business impact)
        chunk_totals = SlideUtils.calculate_column_totals(chunk, 8, "Page Total")
        
        # Create table (8 columns total including business impact)
        data_with_totals = chunk + [chunk_totals]
        table = SlideUtils.create_table_with_headers(slide, len(data_with_totals) + 1, 
                                                   8,  # 8 columns total
                                                   layout['table_left'], Inches(1.2), 
                                                   layout['max_table_width'], Inches(4.8))
        
        # UPDATED COLUMN WIDTHS - Including Business Impact column
        column_widths = [
            Inches(2.0),   # Asset & Network Location (16%)
            Inches(2.7),   # Vulnerability Description (16%)
            Inches(0.7),   # Critical (4%)
            Inches(0.7),   # High (4%)
            Inches(0.8),   # Medium (4%)
            Inches(0.7),   # Low (4%)
            Inches(2.5),   # CVE IDs & CVSS Score (16%)
            Inches(2)    # Business Impact Context (20%) - ADDED BACK
        ]
        
        # Use safe column width setting
        _set_column_widths_safe(table, column_widths)
        
        SlideUtils.set_table_row_heights(table, Inches(0.35), Inches(0.32))
        SlideUtils.format_header_row(table, slide7_data["table"]["columns"])
        SlideUtils.populate_table_data(table, data_with_totals)
        
        # Add footer with navigation info
        if len(patching_chunks) > 1:
            footer_text = f"Critical Vulnerability Response - Page {chunk_idx + 1} of {len(patching_chunks)}"
            SlideUtils.create_footnote(slide, layout['table_left'], Inches(6.2), 
                                     layout['max_table_width'], footer_text, font_size=10)
        
        slides_created.append(slide)
    
    return slides_created


def _set_column_widths_safe(table, widths):
    """Safely set column widths ensuring they are integers"""
    for i, width in enumerate(widths):
        if i < len(table.columns):
            # Convert to int to avoid TypeError
            if hasattr(width, 'value'):  # Handle Inches objects
                width_value = int(width)
            else:
                width_value = int(width)
            table.columns[i].width = width_value


def _create_patching_summary_slide(prs, slide7_data, patching_data):
    """Create a summary overview slide when multiple detail slides exist"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Create title
    SlideUtils.create_title_bar(slide, prs, f"{slide7_data['title']} - Executive Summary")
    
    # Get layout parameters
    layout = SlideUtils.get_standard_layout_params(prs)
    
    # Calculate overall totals (8 columns including business impact)
    overall_totals = SlideUtils.calculate_column_totals(patching_data, 8, "Grand Total")
    
    # Create summary statistics table
    summary_data = [
        ["Metric", "Value"],
        ["Total Critical Assets", str(len(patching_data))],
        ["Critical Vulnerabilities", str(overall_totals[2]) if len(overall_totals) > 2 else "0"],
        ["High Priority Vulnerabilities", str(overall_totals[3]) if len(overall_totals) > 3 else "0"],
        ["Medium Priority Vulnerabilities", str(overall_totals[4]) if len(overall_totals) > 4 else "0"],
        ["Low Priority Vulnerabilities", str(overall_totals[5]) if len(overall_totals) > 5 else "0"],
        ["CVE IDs Identified", "Multiple - See Details"],
        ["Business Impact Areas", str(len(set(row[7] for row in patching_data if len(row) > 7 and row[7])))]  # Count unique business impacts
    ]
    
    # Create summary table
    table = SlideUtils.create_table_with_headers(slide, len(summary_data), 2, 
                                               layout['table_left'], Inches(1.2), 
                                               Inches(8.0), Inches(2.8))
    
    column_widths = [Inches(4.0), Inches(4.0)]
    SlideUtils.set_table_column_widths(table, column_widths)
    SlideUtils.set_table_row_heights(table, Inches(0.35), Inches(0.35))
    SlideUtils.format_header_row(table, summary_data[0])
    SlideUtils.populate_table_data(table, summary_data[1:])
    
    # Add key findings text box
    findings_text = f"""Key Findings - Critical Vulnerability Response:
• {len(patching_data)} critical assets requiring immediate remediation
• Multiple CVE IDs identified with varying CVSS scores
• Business impact assessment completed for each vulnerability
• Asset-specific network location and service context provided
• Detailed breakdown available in following slides for operational teams"""
    
    findings_box = slide.shapes.add_textbox(layout['table_left'], Inches(4.2), 
                                          layout['max_table_width'], Inches(1.8))
    findings_tf = findings_box.text_frame
    findings_tf.text = findings_text
    for paragraph in findings_tf.paragraphs:
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
    
    return slide


def _split_rows_for_slides(rows, max_rows_per_slide):
    """Split list of rows into smaller chunks for pagination"""
    for i in range(0, len(rows), max_rows_per_slide):
        yield rows[i:i + max_rows_per_slide]


def _print_slide7_results(start_time, patching_data, slides_created):
    """Print enhanced slide 7 results with pagination info"""
    runtime = time.time() - start_time
    print(f"Slide 7 series created successfully!")
    print(f"Critical vulnerability items: {len(patching_data)} assets")
    print(f"Total slides created: {slides_created}")
    
    if len(patching_data) > 10:
        print(f"✓ Critical vulnerability data split across multiple slides for readability")
        
    print(f"Runtime: {runtime:.4f} seconds")