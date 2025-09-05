import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from .slide_utils import SlideUtils
from config_colors import COLORS, FONT_SIZES

def create_slide9(prs: Presentation, slide9_data):
    start_time = time.time()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Create title bar
    SlideUtils.create_title_bar(slide, prs, slide9_data["title"])
    
    # Calculate totals
    part1_total, part2_total = _calculate_part_totals(slide9_data)
    grand_total = part1_total + part2_total
    
    # Create side-by-side tables
    _create_eol_tables(slide, slide9_data, part1_total, part2_total, prs)
    
    # Create summary table
    _create_summary_table(slide, part1_total, part2_total, grand_total, prs)
    
    # Print results
    _print_slide9_results(start_time, part1_total, part2_total, grand_total)
    return time.time() - start_time

def _calculate_part_totals(slide9_data):
    """Calculate totals for both parts"""
    part1_total = sum(int(row[1]) for row in slide9_data["tables"][0]["rows"] 
                     if "Total" not in row[0] and row[1].strip().isdigit())
    part2_total = sum(int(row[1]) for row in slide9_data["tables"][1]["rows"] 
                     if "Total" not in row[0] and row[1].strip().isdigit())
    return part1_total, part2_total

def _create_eol_tables(slide, slide9_data, part1_total, part2_total, prs):
    """Create EOL tables side by side"""
    mid_point = prs.slide_width // 2
    table_width = mid_point - Inches(0.6)
    
    # Left table (Part 1)
    _create_eol_table(slide, slide9_data["tables"][0], part1_total, "Part 1 Total", 
                     Inches(0.3), table_width, "table1")
    
    # Right table (Part 2)
    _create_eol_table(slide, slide9_data["tables"][1], part2_total, "Part 2 Total",
                     mid_point + Inches(0.3), table_width, "table2")

def _create_eol_table(slide, table_data, total, total_label, left, width, table_id):
    """Create individual EOL table with totals"""
    filtered_data = [row for row in table_data["rows"] if "Total" not in row[0]]
    data_with_total = filtered_data + [[total_label, str(total)]]
    
    table = SlideUtils.create_table_with_headers(slide, len(data_with_total) + 1, 2,
                                               left, Inches(1), width, Inches(4.5))
    
    SlideUtils.set_table_column_widths(table, [Inches(3.5), Inches(1.5)])
    SlideUtils.set_table_row_heights(table, Inches(0.4), Inches(0.35))
    SlideUtils.format_header_row(table, ["EOL", "Immediate"])
    
    # Custom population for EOL tables with special formatting
    _populate_eol_table_data(table, data_with_total)

def _populate_eol_table_data(table, data_rows):
    """Populate EOL table with special blue formatting for totals"""
    for row_idx, row_data in enumerate(data_rows):
        table_row_idx = row_idx + 1
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(table_row_idx, col_idx)
            is_total = "Total" in str(row_data[0])
            
            if is_total:
                # Blue background for total rows
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*COLORS["blue"])
                cell.text = str(cell_data)
                p = cell.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.color.rgb = RGBColor(*COLORS["white"])
                p.font.size = Pt(FONT_SIZES["table_data"])
                if col_idx == 1:  # Numbers
                    p.alignment = 1
            else:
                # Regular data formatting
                SlideUtils.format_data_cell(cell, cell_data, col_idx, row_idx, False)

def _create_summary_table(slide, part1_total, part2_total, grand_total, prs):
    """Create summary totals table"""
    summary_data = [
        ["Part 1 Total", str(part1_total)],
        ["Part 2 Total", str(part2_total)],
        ["Grand Total", str(grand_total)]
    ]
    
    table = SlideUtils.create_table_with_headers(slide, 4, 2, Inches(0.3), Inches(5.8), 
                                               prs.slide_width // 2 - Inches(0.6), Inches(1.2))
    
    SlideUtils.set_table_column_widths(table, [Inches(3.5), Inches(1.5)])
    SlideUtils.set_table_row_heights(table, Inches(0.4), Inches(0.35))
    SlideUtils.format_header_row(table, ["EOL", "Immediate"])
    
    # Custom formatting for summary table
    _format_summary_table(table, summary_data)

def _format_summary_table(table, summary_data):
    """Apply custom formatting to summary table"""
    colors = [COLORS["very_light_gray"], COLORS["very_light_gray"], COLORS["blue"]]
    font_colors = [COLORS["black"], COLORS["black"], COLORS["white"]]
    bold_flags = [False, False, True]
    
    for row_idx, (row_data, color, font_color, bold) in enumerate(zip(summary_data, colors, font_colors, bold_flags)):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*color)
            
            p = cell.text_frame.paragraphs[0]
            p.font.bold = bold
            p.font.size = Pt(FONT_SIZES["table_data"])
            p.font.color.rgb = RGBColor(*font_color)
            if col_idx == 1:  # Numbers
                p.alignment = 1

def _print_slide9_results(start_time, part1_total, part2_total, grand_total):
    """Print slide 9 results"""
    runtime = time.time() - start_time
    print(f"Slide 9 created successfully!")
    print(f"Calculated Totals - Part 1: {part1_total}, Part 2: {part2_total}, Grand Total: {grand_total}")
    print(f"Runtime: {runtime:.4f} seconds")
