from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from config_colors import COLORS, FONT_SIZES

class SlideUtils:
    @staticmethod
    def create_title_bar(slide, prs, title_text):
        """Create standardized blue title bar for all slides"""
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = Inches(0.6)
        
        # Create title shape
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS["blue"])
        shape.line.color.rgb = RGBColor(*COLORS["blue"])
        
        # Add title text
        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(FONT_SIZES["title"])
        p.font.color.rgb = RGBColor(*COLORS["white"])
        p.alignment = 1
        
        return title_box
    
    @staticmethod
    def create_subtitle(slide, left, top, width, subtitle_text, font_size=16):
        """Create standardized subtitle"""
        subtitle_box = slide.shapes.add_textbox(left, top, width, Inches(0.3))
        subtitle_tf = subtitle_box.text_frame
        subtitle_p = subtitle_tf.paragraphs[0]
        subtitle_p.text = subtitle_text
        subtitle_p.font.bold = True
        subtitle_p.font.size = Pt(font_size)
        subtitle_p.font.color.rgb = RGBColor(*COLORS["black"])
        subtitle_p.alignment = 1
        return subtitle_box
    
    @staticmethod
    def create_table_with_headers(slide, rows, cols, left, top, width, height):
        """Create table with standard dimensions and return table object"""
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        return table
    
    @staticmethod
    def set_table_column_widths(table, widths):
        """Set column widths for table"""
        for i, width in enumerate(widths):
            if i < len(table.columns):
                table.columns[i].width = width
    
    @staticmethod
    def set_table_row_heights(table, header_height=Inches(0.4), data_height=Inches(0.3)):
        """Set row heights for table"""
        if len(table.rows) > 0:
            table.rows[0].height = header_height  # Header row
        for i in range(1, len(table.rows)):
            table.rows[i].height = data_height  # Data rows
    
    @staticmethod
    def format_header_row(table, columns, row_idx=0):
        """Format table header row with blue background"""
        for i, col_name in enumerate(columns):
            if i < len(table.columns):
                cell = table.cell(row_idx, i)
                cell.text = col_name
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*COLORS["blue"])
                
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.bold = True
                paragraph.font.size = Pt(FONT_SIZES["table_header"])
                paragraph.font.color.rgb = RGBColor(*COLORS["white"])
                paragraph.alignment = 1
    
    @staticmethod
    def format_data_cell(cell, cell_data, col_idx, row_idx, is_total_row=False):
        """Format individual data cell with consistent styling"""
        cell.text = str(cell_data) if cell_data else ""
        paragraph = cell.text_frame.paragraphs[0]
        
        if is_total_row:
            # Blue background for total rows
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*COLORS["blue"])
            paragraph.font.bold = True
            paragraph.font.size = Pt(FONT_SIZES["table_totals"])
            paragraph.font.color.rgb = RGBColor(*COLORS["white"])
        else:
            # Regular data cell formatting
            paragraph.font.size = Pt(FONT_SIZES["table_data"])
            paragraph.font.color.rgb = RGBColor(*COLORS["black"])
            
            # Alternate row coloring
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*COLORS["very_light_gray"])
        
        # Center align numbers (columns > 0)
        if col_idx > 0:
            paragraph.alignment = 1
    
    @staticmethod
    def populate_table_data(table, data_rows, start_row=1):
        """Populate table with data rows"""
        for row_idx, row_data in enumerate(data_rows):
            table_row_idx = start_row + row_idx
            if table_row_idx < len(table.rows):
                for col_idx, cell_data in enumerate(row_data):
                    if col_idx < len(table.columns):
                        cell = table.cell(table_row_idx, col_idx)
                        is_total = isinstance(row_data, list) and len(row_data) > 0 and "Total" in str(row_data[0])
                        SlideUtils.format_data_cell(cell, cell_data, col_idx, row_idx, is_total)
    
    @staticmethod
    def calculate_column_totals(data_rows, num_cols, label="Grand Total"):
        """Calculate totals for numeric columns"""
        totals = [label] + [0] * (num_cols - 2) + [0]  # label + numeric cols + grand total
        
        for row in data_rows:
            if "Total" not in str(row[0]):  # Skip existing total rows
                for col_idx in range(1, num_cols - 1):  # Skip label and grand total columns
                    value = str(row[col_idx]).strip() if len(row) > col_idx else "0"
                    if value.isdigit():
                        totals[col_idx] += int(value)
        
        # Calculate grand total
        totals[-1] = sum(totals[1:-1])
        return totals
    
    @staticmethod
    def create_footnote(slide, left, top, width, footnote_text, font_size=8):
        """Create standardized footnote"""
        footnote_box = slide.shapes.add_textbox(left, top, width, Inches(0.2))
        footnote_tf = footnote_box.text_frame
        footnote_p = footnote_tf.paragraphs[0]
        footnote_p.text = footnote_text
        footnote_p.font.size = Pt(font_size)
        footnote_p.font.italic = True
        footnote_p.font.color.rgb = RGBColor(*COLORS["black"])
        return footnote_box

    @staticmethod
    def get_standard_layout_params(prs):
        """Get standard layout parameters"""
        return {
            'max_table_width': prs.slide_width - Inches(1.0),
            'table_left': Inches(0.5),
            'standard_gap': Inches(0.4),
            'subtitle_height': Inches(0.3)
        }

    @staticmethod
    def create_two_column_layout(prs):
        """Get parameters for two-column layout"""
        mid_point = prs.slide_width // 2
        return {
            'left_column_start': Inches(0.3),
            'right_column_start': mid_point + Inches(0.3),
            'column_width': mid_point - Inches(0.6),
            'mid_point': mid_point
        }

    @staticmethod
    def apply_severity_color_coding(cell, severity):
        """Apply color coding based on severity level"""
        severity_colors = {
            "Immediate": (255, 0, 0),
            "Critical": (220, 20, 60),
            "High": (255, 165, 0),
            "Medium": (255, 255, 0),
            "Low": (128, 128, 128)
        }
        
        if severity in severity_colors:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*severity_colors[severity])
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.color.rgb = RGBColor(*COLORS["white"])

    @staticmethod
    def filter_total_rows(data_rows):
        """Filter out existing total rows from data"""
        return [row for row in data_rows if "Total" not in str(row[0])]

    @staticmethod
    def create_legends_section(slide, left, top, width, legends_data):
        """Create legends section with formatting"""
        legends_box = slide.shapes.add_textbox(left, top, width, Inches(1.0))
        legends_tf = legends_box.text_frame
        
        # Title
        legends_title = legends_tf.paragraphs[0]
        legends_title.text = "Legends"
        legends_title.font.bold = True
        legends_title.font.size = Pt(12)
        
        # Legend items
        for legend in legends_data:
            legend_p = legends_tf.add_paragraph()
            legend_p.text = f"• {legend['term']}: {legend['description']}"
            legend_p.font.size = Pt(8)
            legend_p.space_before = Pt(3)
        
        return legends_box
