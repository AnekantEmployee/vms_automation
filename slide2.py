from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from ppt_data import TEXT_BLACK, LIGHT_BLUE, DARK_BLUE, WHITE
from ppt_data import CHART_CATEGORIES, CHART_VALUES, TABLE_DATA, RECOMMENDATIONS_DATA, SECOND_TABLE_DATA
 
def create_header(slide, text, left=Cm(1), top=Cm(1), width=Cm(6), height=Cm(0.95),
                 font_size=Cm(0.5), bg_color=DARK_BLUE, text_color=WHITE):
    """
    Create a styled header textbox on a slide
    """
    header = slide.shapes.add_textbox(left, top, width, height)
   
    # Style the header
    header_frame = header.text_frame
    header_frame.text = text
    header_para = header_frame.paragraphs[0]
    header_para.alignment = PP_ALIGN.LEFT
   
    # Style header text
    header_run = header_para.runs[0]
    header_run.font.size = font_size
    header_run.font.bold = True
    header_run.font.color.rgb = text_color
   
    # Add colored background to header
    header.fill.solid()
    header.fill.fore_color.rgb = bg_color
   
    return header
 
def create_yash_recommendations_section(slide, RECOMMENDATIONS_DATA):
    """
    Create the "YASH Recommendations - Why Learn More?" section (40% width on right side)
    with proper blue circular icons
    """
    # Fixed positioning to match the "What's in it for you?" section
    recommendations_left = Cm(12.5)  # Position on right side
    recommendations_top = Cm(17)     # ALIGNED WITH "What's in it for you?" SECTION (changed from Cm(16))
    recommendations_width = Cm(8.5)  # 40% of slide width
    recommendations_height = Cm(10)  # Appropriate height
 
    recommendations_box = slide.shapes.add_textbox(
        recommendations_left, recommendations_top, recommendations_width, recommendations_height
    )
 
    # Style the text box
    recommendations_frame = recommendations_box.text_frame
    recommendations_frame.word_wrap = True
 
    # Add main title
    title_para = recommendations_frame.paragraphs[0]
    title_run = title_para.add_run()
    title_run.text = "YASH Recommendations"
    title_run.font.size = Cm(0.6)
    title_run.font.bold = True
    title_run.font.color.rgb = TEXT_BLACK
    title_para.alignment = PP_ALIGN.LEFT
 
    # Add subtitle
    subtitle_para = recommendations_frame.add_paragraph()
    subtitle_run = subtitle_para.add_run()
    subtitle_run.text = "Why Learn More?"
    subtitle_run.font.size = Cm(0.5)
    subtitle_run.font.color.rgb = TEXT_BLACK
    subtitle_para.alignment = PP_ALIGN.LEFT
 
    # Add spacing
    recommendations_frame.add_paragraph()
 
    # Create circular blue backgrounds for icons and text
    current_y = recommendations_top + Cm(2)  # Start position for first item
 
    for i, rec in enumerate(RECOMMENDATIONS_DATA):
        # Create blue circle for icon
        circle_left = recommendations_left + Cm(0.2)
        circle_top = current_y
        circle_size = Cm(1.2)
 
        # Add blue circle shape
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            circle_left, circle_top, circle_size, circle_size
        )
 
        # Style the circle
        circle.fill.solid()
        circle.fill.fore_color.rgb = LIGHT_BLUE
        circle.line.fill.background()  # Remove border
 
        # Add icon text in center of circle
        icon_box = slide.shapes.add_textbox(
            circle_left, circle_top, circle_size, circle_size
        )
 
        icon_frame = icon_box.text_frame
        icon_frame.margin_left = 0
        icon_frame.margin_right = 0
        icon_frame.margin_top = Cm(0.25)
        icon_frame.margin_bottom = 0
        icon_para = icon_frame.paragraphs[0]
        icon_para.alignment = PP_ALIGN.CENTER
        icon_run = icon_para.add_run()
        icon_run.text = rec['icon']
        icon_run.font.size = Cm(0.5)
        icon_run.font.bold = True  # Make icon bold
 
        # Add title and subtitle text next to circle
        text_left = circle_left + circle_size + Cm(0.5)
        text_top = circle_top
        text_width = Cm(6)
        text_height = Cm(1.5)
 
        text_box = slide.shapes.add_textbox(
            text_left, text_top, text_width, text_height
        )
 
        text_frame = text_box.text_frame
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = Cm(0.1)
        text_frame.margin_bottom = 0
 
        # Add title
        title_para = text_frame.paragraphs[0]
        title_run = title_para.add_run()
        title_run.text = rec['title']
        title_run.font.size = Cm(0.42)
        title_run.font.bold = True
        title_run.font.color.rgb = TEXT_BLACK
        title_para.alignment = PP_ALIGN.LEFT

        # Add subtitle (also bold and same size as title)
        subtitle_para = text_frame.add_paragraph()
        subtitle_run = subtitle_para.add_run()
        subtitle_run.text = rec['subtitle']
        subtitle_run.font.size = Cm(0.42)   # Same size as title
        subtitle_run.font.bold = True      # Make bold
        subtitle_run.font.color.rgb = TEXT_BLACK
        subtitle_para.alignment = PP_ALIGN.LEFT
 
        # Make text background transparent
        text_box.fill.background()
 
        # Move to next position
        current_y += Cm(1.75)
 
    # Add framework note at the bottom
    framework_top = current_y + Cm(0)
    framework_width = Cm(8)
    framework_height = Cm(1)
 
    framework_box = slide.shapes.add_textbox(
        recommendations_left, framework_top, framework_width, framework_height
    )
 
    framework_frame = framework_box.text_frame
    framework_para = framework_frame.paragraphs[0]
    framework_run = framework_para.add_run()
    framework_run.text = "(e.g. NIST & ISO/IEC 27001)"
    framework_run.font.size = Cm(0.35)
    framework_run.font.color.rgb = TEXT_BLACK
    framework_run.font.italic = True
    framework_para.alignment = PP_ALIGN.LEFT
 
    # Make background transparent
    recommendations_box.fill.background()
    framework_box.fill.background()
 
    return recommendations_box

def create_vulnerability_summary_slide(p, img_path='ppt-generation/bg/2.jpg'):
    """
    Create the vulnerability summary slide with all content
    """
    # Data configuration
    CHART_DATA = {
        'categories': CHART_CATEGORIES,
        'values': CHART_VALUES,
        'colors': [
            RGBColor(180, 0, 0),    # Critical - Dark red
            RGBColor(220, 0, 0),    # High - Red
            RGBColor(255, 165, 0),  # Medium - Orange
            RGBColor(255, 220, 0)   # Low - Yellow
        ]
    }
     
    CHART_CONFIG = {
        'left': Cm(1),
        'top': Cm(2.5),
        'width': Cm(9),
        'height': Cm(8),
        'title': "Vulnerability Count"
    }
     
    TABLE_CONFIG = {
        'rows': 9,
        'cols': 2,
        'left': Cm(11.5),
        'top': Cm(2.5),
        'width': Cm(8.5),
        'height': Cm(8),
        'column_widths': [Cm(6), Cm(2.5)]
    }
    
    # Add slide with blank layout and background
    slide = p.slides.add_slide(p.slide_layouts[6])
    
    # Add background image
    left = top = 0
    pic = slide.shapes.add_picture(img_path, left, top, width=p.slide_width, height=p.slide_height)
    
    # Send picture to back so other elements appear on top
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)
     
    # Add header using the function
    create_header(slide, "Vulnerability Summary")
     
    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = CHART_DATA['categories']
    chart_data.add_series('Vulnerabilities', CHART_DATA['values'])
     
    # Add chart
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        CHART_CONFIG['left'],
        CHART_CONFIG['top'],
        CHART_CONFIG['width'],
        CHART_CONFIG['height'],
        chart_data
    )
     
    # Customize chart
    chart_obj = chart.chart
     
    # Set chart title
    chart_obj.has_title = True
    chart_obj.chart_title.text_frame.text = CHART_CONFIG['title']
    chart_obj.chart_title.text_frame.paragraphs[0].font.size = Cm(0.4)
    chart_obj.chart_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
     
    # Get series and customize colors
    series = chart_obj.series[0]
     
    try:
        points = series.points
        for i, point in enumerate(points):
            if i < len(CHART_DATA['colors']):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = CHART_DATA['colors'][i]
    except Exception:
        # Fallback if individual point coloring fails
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = RGBColor(255, 165, 0)
     
    # Customize axes
    try:
        value_axis = chart_obj.value_axis
        value_axis.maximum_scale = 5.5
        value_axis.minimum_scale = 0
        value_axis.major_unit = 1
       
        # Clean up axis appearance
        value_axis.tick_labels.font.size = Cm(0.3)
        category_axis = chart_obj.category_axis
        category_axis.tick_labels.font.size = Cm(0.3)
       
        # Make gridlines white
        value_axis.major_gridlines.format.line.color.rgb = WHITE
    except Exception:
        pass
     
    # Add data labels
    try:
        series.has_data_labels = True
        data_labels = series.data_labels
        data_labels.show_value = True
        data_labels.font.size = Cm(0.3)
        data_labels.font.color.rgb = RGBColor(64, 64, 64)
    except Exception:
        pass
     
    # Remove legend and clean up appearance
    try:
        chart_obj.has_legend = False
        chart_obj.plot_area.format.line.color.rgb = WHITE
        chart.format.line.color.rgb = WHITE
    except Exception:
        pass
     
    # Add table
    table_shape = slide.shapes.add_table(
        rows=TABLE_CONFIG['rows'],
        cols=TABLE_CONFIG['cols'],
        left=TABLE_CONFIG['left'],
        top=TABLE_CONFIG['top'],
        width=TABLE_CONFIG['width'],
        height=TABLE_CONFIG['height']
    )
     
    table = table_shape.table
     
    # Set column widths
    for i, width in enumerate(TABLE_CONFIG['column_widths']):
        table.columns[i].width = width
     
    # Populate table and style
    for row_idx, row_data in enumerate(TABLE_DATA):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_data
           
            # Style header row
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
               
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                paragraph = cell.text_frame.paragraphs[0]
                run = paragraph.runs[0]
                run.font.color.rgb = WHITE
                run.font.bold = True
                run.font.size = Cm(0.35)
                paragraph.alignment = PP_ALIGN.LEFT
            else:
                # Data row styling
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(230, 240, 250)
               
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertical center alignment
                paragraph = cell.text_frame.paragraphs[0]
                run = paragraph.runs[0]
                run.font.size = Cm(0.3)
                run.font.color.rgb = RGBColor(64, 64, 64)
                paragraph.alignment = PP_ALIGN.LEFT
     
    # Add table borders
    for row in table.rows:
        for cell in row.cells:
            cell.margin_left = Cm(0.1)
            cell.margin_right = Cm(0.1)
            cell.margin_top = Cm(0.05)
            cell.margin_bottom = Cm(0.05)
     
    # Add second table below header - full width

     
    second_table_shape = slide.shapes.add_table(
        rows=4,     # 3 threat types + 1 header
        cols=2,     # Threat Type and Risks columns
        left=Cm(1),      # Full width starting from left
        top=Cm(11),      # Below the first table and chart
        width=Cm(19),    # Full width
        height=Cm(4.5)     # Reduced height for 4 rows
    )
     
    second_table = second_table_shape.table
     
    # Set column widths for second table
    second_table.columns[0].width = Cm(4)   # Threat Type column
    second_table.columns[1].width = Cm(15)  # Risks column (wider for longer text)
     
    # Populate second table and style
    for row_idx, row_data in enumerate(SECOND_TABLE_DATA):
        for col_idx, cell_data in enumerate(row_data):
            cell = second_table.cell(row_idx, col_idx)
            cell.text = cell_data
           
            # Style header row
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
               
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                paragraph = cell.text_frame.paragraphs[0]
                run = paragraph.runs[0]
                run.font.color.rgb = WHITE
                run.font.bold = True
                run.font.size = Cm(0.35)
                paragraph.alignment = PP_ALIGN.LEFT
            else:
                # Data row styling
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(230, 240, 250)  # Light blue alternating
               
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                paragraph = cell.text_frame.paragraphs[0]
                run = paragraph.runs[0]
                run.font.size = Cm(0.3)
                run.font.color.rgb = RGBColor(64, 64, 64)
               
                # Make first column (Threat Type) bold
                if col_idx == 0:
                    run.font.bold = True
               
                paragraph.alignment = PP_ALIGN.LEFT
     
    # Add borders to second table
    for row in second_table.rows:
        for cell in row.cells:
            cell.margin_left = Cm(0.1)
            cell.margin_right = Cm(0.1)
            cell.margin_top = Cm(0.05)
            cell.margin_bottom = Cm(0.05)
     
    # Add "What's in it for you?" section above footer (left 60%)
    whats_in_it_box = slide.shapes.add_textbox(
        Cm(1),      # Left position
        Cm(17),     # Position above footer
        Cm(11),   # 50% width
        Cm(10)      # Height
    )
     
    # Style the text box
    whats_in_it_frame = whats_in_it_box.text_frame
    whats_in_it_frame.word_wrap = True
     
    # Add title
    title_para = whats_in_it_frame.paragraphs[0]
    title_run = title_para.add_run()
    title_run.text = "What's in it for you?"
    title_run.font.size = Cm(0.6)  # Large title
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 0, 0)  # Dark gray
    title_para.alignment = PP_ALIGN.LEFT
     
    # Add bullet points
    bullet_points = [
        "Explore the most prevalent and impactful threats, techniques, and trends that we've observed.",
        "Prioritize and categorize the threats based on their severity and impact on your business and invest your valuable time and resources on stuff that needs immediate attention.",
        "Shape and inform your readiness, detection, and response to critical threats",
        "Get a custom consultation from our cybersecurity experts and secure your IT landscape."
    ]
     
    for bullet_text in bullet_points:
        bullet_para = whats_in_it_frame.add_paragraph()
        bullet_para.space_before = Pt(6)  # Approx 0.15 cm space after bullet paragraph
        bullet_para.level = 0  # Bullet level
        bullet_run = bullet_para.add_run()
        bullet_run.text = f"- {bullet_text}"
        bullet_run.font.size = Cm(0.36)  # Regular text size
        bullet_run.font.color.rgb = RGBColor(6, 6, 6)  # Dark gray
        bullet_para.alignment = PP_ALIGN.LEFT

    contact_para2 = whats_in_it_frame.add_paragraph()
    contact_para2.space_before = Pt(18)  # Adds space before this paragraph

    contact_run1 = contact_para2.add_run()
    contact_run1.text = "For more information write to us at"
    contact_run1.font.size = Cm(0.36)
    contact_run1.font.color.rgb = RGBColor(6, 6, 6)

    contact_para3 = whats_in_it_frame.add_paragraph()
    contact_run2 = contact_para3.add_run()
    contact_run2.text = "cybersecurity@yash.com"
    contact_run2.font.size = Cm(0.4)
    contact_run2.font.bold = True
    contact_run2.font.color.rgb = RGBColor(0, 0, 0)
    contact_para3.alignment = PP_ALIGN.LEFT
     
    # Make background transparent
    whats_in_it_box.fill.background()
     
    # Add YASH Recommendations section to the right of the "What's in it for you?" section
    create_yash_recommendations_section(slide, RECOMMENDATIONS_DATA)
     
    # Add disclaimer footer at the bottom - full width blue strip
    footer_full = slide.shapes.add_textbox(
        Cm(0),      # No left padding - start from edge
        Cm(26),   # Bottom position
        Cm(21),     # Full width - no right padding
        Cm(2)       # Height
    )
     
    # Add blue background to full strip
    footer_full.fill.solid()
    footer_full.fill.fore_color.rgb = LIGHT_BLUE
     
    # Add disclaimer text (left 50%)
    disclaimer_left = slide.shapes.add_textbox(
        Cm(0.5),    # Small margin from edge
        Cm(26.4),   # Slightly adjusted position
        Cm(10.5),   # 50% width
        Cm(1.6)     # Height
    )
     
    # Style disclaimer text
    disclaimer_frame = disclaimer_left.text_frame
    disclaimer_frame.word_wrap = True  # Enable text wrapping
    disclaimer_frame.margin_left = 0
    disclaimer_frame.margin_right = 0
    disclaimer_frame.margin_top = 0
    disclaimer_frame.margin_bottom = 0
     
    # Add disclaimer text with red "Disclaimer:"
    disclaimer_para = disclaimer_frame.paragraphs[0]
    disclaimer_run1 = disclaimer_para.add_run()  # Create first run
    disclaimer_run1.text = "Disclaimer:"
    disclaimer_run1.font.size = Cm(0.28)  # 8pt
    disclaimer_run1.font.color.rgb = RGBColor(255, 0, 0)  # Red
     
    # Add rest of text in white
    disclaimer_run2 = disclaimer_para.add_run()
    disclaimer_run2.text = " All the information for the report has been obtained from publicly available sources. YASH technologies does not perform any unauthorized assessment that could impact your business"
    disclaimer_run2.font.size = Cm(0.28)  # 8pt
    disclaimer_run2.font.color.rgb = WHITE  # White
     
    disclaimer_para.alignment = PP_ALIGN.LEFT
     
    # Add tagline text (right 50%)
    tagline_right = slide.shapes.add_textbox(
        Cm(11.5),   # Start from 50% width
        Cm(26.4),   # Same position as disclaimer
        Cm(10),     # Right 50% width
        Cm(1.6)     # Height
    )
     
    # Style tagline text
    tagline_frame = tagline_right.text_frame
    tagline_frame.margin_left = 0
    tagline_frame.margin_right = 0
    tagline_frame.margin_top = 0
    tagline_frame.margin_bottom = 0
     
    tagline_frame.text = "Maximise security posture with"
    tagline_para = tagline_frame.paragraphs[0]
    tagline_run = tagline_para.runs[0]
    tagline_run.font.size = Cm(0.35)  # Reduced font size (was 0.42)
    tagline_run.font.color.rgb = WHITE  # White text
    tagline_run.font.bold = True
     
    # Add second line with white color for SOC/MDR/VMS/TRPM
    tagline_para2 = tagline_frame.add_paragraph()
    tagline_run2 = tagline_para2.add_run()
    tagline_run2.text = "SOC/MDR/VMS/TRPM"
    tagline_run2.font.size = Cm(0.42)  # Same reduced font size
    tagline_run2.font.color.rgb = WHITE  # White text (changed from black)
    tagline_run2.font.bold = True
    tagline_para2.alignment = PP_ALIGN.LEFT
     
    # Make background transparent for text boxes (they sit on top of blue strip)
    disclaimer_left.fill.background()
    tagline_right.fill.background()
    
    return slide
