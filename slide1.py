from pptx.oxml import parse_xml
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from ppt_data import DARK_BLUE, LIGHT_SKY_BLUE, YELLOW, VULNERABILITIES


def add_vulnerability_content(slide, vulnerabilities, title="Cyberattack Entry Points Identified", content_cfg=""):
    # Calculate total content height needed
    item_height = content_cfg['item_height']
    items_spacing = content_cfg['items_spacing']
    total_content_height = (
        Inches(0.7) +  # Title height with padding
        (len(vulnerabilities) * item_height) + 
        ((len(vulnerabilities) - 1) * items_spacing)
    )
    
    # Calculate starting Y position to center everything vertically
    available_height = content_cfg.get('available_height', Inches(5.5))  # Default slide content height
    start_y = content_cfg['y'] + (available_height - total_content_height) / 2
    
    # Add title background
    title_left = content_cfg['x']
    title_top = start_y
    title_width = content_cfg['width']
    title_height = Inches(0.35)
    
    # Add background rectangle for title
    title_bg = slide.shapes.add_shape(
        1,  # Rectangle shape
        title_left,
        title_top,
        title_width,
        title_height + Inches(0.2)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = DARK_BLUE
    title_bg.line.fill.background()
    
    # Add textbox with vertical centering
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height + Inches(0.2))
    title_frame = title_box.text_frame
    title_frame.text = title
    
    # Format title with vertical centering
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_run = title_p.runs[0]
    title_run.font.name = 'Arial'
    title_run.font.size = Pt(20)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(255, 255, 255)  # White text
    
    # Move title background behind text
    slide.shapes._spTree.remove(title_bg._element)
    slide.shapes._spTree.insert(-2, title_bg._element)
    
    # Calculate content area dimensions
    items_start_y = title_top + title_height + Inches(0.25)
    
    # Calculate text positioning
    text_left = content_cfg['x'] + content_cfg['text_left_offset']
    text_width = content_cfg['width'] * content_cfg['text_width_ratio']
    
    # Add vulnerability items
    for i, vuln in enumerate(vulnerabilities):
        y_position = items_start_y + (i * (item_height + items_spacing))
        
        # Add background for each item
        item_bg = slide.shapes.add_shape(
            1,  # Rectangle shape
            content_cfg['x'],
            y_position,
            content_cfg['width'],
            item_height
        )
        item_bg.fill.solid()
        item_bg.fill.fore_color.rgb = LIGHT_SKY_BLUE  # Light blue background
        item_bg.line.color.rgb = RGBColor(200, 220, 240)
        item_bg.line.width = Pt(1)
        
        # Add icon
        if vuln.get('icon'):
            icon_left = content_cfg['x'] + content_cfg['icon_margin']
            icon_top = y_position + (item_height - content_cfg['icon_size']) / 2
            icon = slide.shapes.add_picture(
                vuln['icon'],
                icon_left,
                icon_top,
                width=content_cfg['icon_size'],
                height=content_cfg['icon_size']
            )
        
        # Add text with vertical centering
        text_box = slide.shapes.add_textbox(
            text_left,
            y_position,
            text_width,
            item_height
        )
        text_frame = text_box.text_frame
        text_frame.text = vuln['text']
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertical centering
        
        # Format text
        text_p = text_frame.paragraphs[0]
        text_p.alignment = PP_ALIGN.LEFT
        text_run = text_p.runs[0]
        text_run.font.name = 'Arial'
        text_run.font.size = Pt(13)
        text_run.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

def add_yellow_box_custom_corners(slide):
    """
    Create a yellow box with only top-right corner rounded using XML manipulation
    
    Args:
        slide: PowerPoint slide object
        x, y: Position
        width, height: Dimensions  
        top_right_radius: Corner radius in pixels for top-right corner
        border_width: Border thickness
    """
    x=Cm(0.5)
    y=Cm(11.15)
    width=Cm(8)
    height=Cm(3.5)
    top_right_radius=22
    border_width=Inches(0.02)
    
    # Convert measurements to EMU (English Metric Units)
    x_emu = x.emu
    y_emu = y.emu
    width_emu = width.emu
    height_emu = height.emu
    radius_emu = top_right_radius * 9525  # Convert pixels to EMU
    
    # Create custom shape XML with individual corner radii
    shape_xml = f'''
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <p:nvSpPr>
            <p:cNvPr id="1" name="CustomYellowBox"/>
            <p:cNvSpPr/>
            <p:nvPr/>
        </p:nvSpPr>
        <p:spPr>
            <a:xfrm>
                <a:off x="{x_emu}" y="{y_emu}"/>
                <a:ext cx="{width_emu}" cy="{height_emu}"/>
            </a:xfrm>
            <a:custGeom>
                <a:avLst/>
                <a:gdLst/>
                <a:ahLst/>
                <a:cxnLst/>
                <a:rect l="0" t="0" r="{width_emu}" b="{height_emu}"/>
                <a:pathLst>
                    <a:path w="{width_emu}" h="{height_emu}">
                        <a:moveTo>
                            <a:pt x="0" y="0"/>
                        </a:moveTo>
                        <a:lnTo>
                            <a:pt x="{width_emu - radius_emu}" y="0"/>
                        </a:lnTo>
                        <a:arcTo wR="{radius_emu}" hR="{radius_emu}" stAng="16200000" swAng="5400000"/>
                        <a:lnTo>
                            <a:pt x="{width_emu}" y="{height_emu}"/>
                        </a:lnTo>
                        <a:lnTo>
                            <a:pt x="0" y="{height_emu}"/>
                        </a:lnTo>
                        <a:close/>
                    </a:path>
                </a:pathLst>
            </a:custGeom>
            <a:solidFill>
                <a:srgbClr val="FFC000"/>
            </a:solidFill>
            <a:ln w="{int(border_width.emu)}">
                <a:solidFill>
                    <a:srgbClr val="FFC000"/>
                </a:solidFill>
            </a:ln>
            <a:effectLst>
                <a:outerShdw blurRad="76200" dist="38100" dir="2700000" rotWithShape="0">
                    <a:srgbClr val="000000">
                        <a:alpha val="40000"/>
                    </a:srgbClr>
                </a:outerShdw>
            </a:effectLst>
        </p:spPr>
        <p:txBody>
            <a:bodyPr/>
            <a:lstStyle/>
            <a:p/>
        </p:txBody>
    </p:sp>'''
    
    # Parse and add the custom shape
    shape_element = parse_xml(shape_xml)
    slide.shapes._spTree.append(shape_element)
    
    # Create main textbox
    textbox = slide.shapes.add_textbox(x, y, width, height)
    textbox.fill.background()  # Transparent fill
    textbox.line.fill.background()  # Transparent border

    # Access the text frame
    text_frame = textbox.text_frame
    text_frame.margin_left = Cm(0.3)  # Small margin
    text_frame.margin_right = Cm(0.2)
    text_frame.margin_top = Cm(0.1)
    text_frame.margin_bottom = Cm(0.1)
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertical center
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.word_wrap = False  # Disable word wrap for table-like layout

    # Clear default paragraph
    text_frame.clear()

    # Calculate half width for positioning
    half_width = width / 2

    # Method 1: Using two separate textboxes (Recommended)
    # Left textbox (50% width)
    left_textbox = slide.shapes.add_textbox(x, y, half_width, height)
    left_textbox.fill.background()  # Transparent fill
    left_textbox.line.fill.background()  # Transparent border

    left_text_frame = left_textbox.text_frame
    left_text_frame.margin_left = Cm(0.3)
    left_text_frame.margin_right = Cm(0.2)
    left_text_frame.margin_top = Cm(0.1)
    left_text_frame.margin_bottom = Cm(0.1)
    left_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    left_text_frame.auto_size = MSO_AUTO_SIZE.NONE
    left_text_frame.word_wrap = True

    # Clear and add content to left textbox
    left_text_frame.clear()
    left_paragraph = left_text_frame.paragraphs[0]
    left_paragraph.alignment = PP_ALIGN.LEFT  # Changed to LEFT alignment

    # Add "Your Risk Score" text
    left_title_run = left_paragraph.add_run()
    left_title_run.text = "Your Risk Score"
    left_title_run.font.size = Pt(20)
    left_title_run.font.color.rgb = RGBColor(0, 0, 0)  # Black text

    # Add line break
    left_paragraph.add_run().text = "\n"

    # Add "63" text
    left_score_run = left_paragraph.add_run()
    left_score_run.text = "63"
    left_score_run.font.size = Pt(38)  # Larger font for the score
    left_score_run.font.bold = True
    left_score_run.font.color.rgb = RGBColor(0, 0, 0)  # Black text

    # Right textbox (50% width)
    right_textbox = slide.shapes.add_textbox(x + half_width, y, half_width, height)
    right_textbox.fill.background()  # Transparent fill
    right_textbox.line.fill.background()  # Transparent border

    right_text_frame = right_textbox.text_frame
    right_text_frame.margin_left = Cm(0.3)
    right_text_frame.margin_right = Cm(0.2)
    right_text_frame.margin_top = Cm(0.1)
    right_text_frame.margin_bottom = Cm(0.1)
    right_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    right_text_frame.auto_size = MSO_AUTO_SIZE.NONE
    right_text_frame.word_wrap = True

    # Clear and add content to right textbox
    right_text_frame.clear()
    right_paragraph = right_text_frame.paragraphs[0]
    right_paragraph.alignment = PP_ALIGN.LEFT  # Changed to LEFT alignment

    # Add right content (example)
    right_title_run = right_paragraph.add_run()
    right_title_run.text = "Industry Average"
    right_title_run.font.size = Pt(20)
    right_title_run.font.color.rgb = RGBColor(0, 0, 0)  # Black text

    # Add line break
    right_paragraph.add_run().text = "\n"

    # Add right score
    right_score_run = right_paragraph.add_run()
    right_score_run.text = "45"
    right_score_run.font.size = Pt(38)  # Larger font for the score
    right_score_run.font.bold = True
    right_score_run.font.color.rgb = RGBColor(0, 0, 0)  # Black text

    # Optional: Group the textboxes together
    # Note: You'll need to remove the original textbox if you're using this method
    try:
        # Remove the original textbox if it was created
        sp = slide.shapes._spTree
        sp.remove(textbox._element)
    except:
        pass
    
    return slide.shapes[-1]

def add_transparent_box_with_text(slide, heading, text, x, y, width, height):
    """
    Add a transparent box with heading and text content
    
    Args:
        slide: PowerPoint slide object
        heading: Header text string
        text: Main content text string  
        x: X position (use Cm() or Inches())
        y: Y position (use Cm() or Inches())
        width: Box width (use Cm() or Inches())
        height: Box height (use Cm() or Inches())
    """
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Cm, Inches
    
    # Add transparent rectangle shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, width, height
    )
    
    # Make shape completely transparent by removing fill
    fill = shape.fill
    fill.background()  # This removes the fill entirely
    
    # Remove the outline/border
    line = shape.line
    line.fill.background()

    # Add heading text box inside the transparent box
    txBox_heading = slide.shapes.add_textbox(
        x, y, width, Cm(1.2)
    )
    tf_heading = txBox_heading.text_frame
    p_heading = tf_heading.paragraphs[0]
    run_heading = p_heading.add_run()
    run_heading.text = heading
    run_heading.font.bold = True
    run_heading.font.size = Inches(0.3)  # Approximately 24pt
    run_heading.font.color.rgb = RGBColor(0, 0, 0)  # Black text

    # Add main text box inside the transparent box
    txBox_text = slide.shapes.add_textbox(
        x, y + Cm(1), width, height
    )
    tf_text = txBox_text.text_frame
    tf_text.word_wrap = True  # Enable word wrapping
    p_text = tf_text.paragraphs[0]
    run_text = p_text.add_run()
    run_text.text = text
    run_text.font.size = Inches(0.2)  # Approximately 16pt
    run_text.font.color.rgb = RGBColor(0, 0, 0)  # Black text
    
    return shape

def create_blank_slide_with_background(p, img_path, top_right_img_path='ppt-generation/bg/yash-logo.png', margin=40):
    # Add slide with blank layout
    slide = p.slides.add_slide(p.slide_layouts[6])
    
    # Add background image
    left = top = 0
    bg_pic = slide.shapes.add_picture(img_path, left, top, width=p.slide_width, height=p.slide_height)
    
    # Send background picture to back
    slide.shapes._spTree.remove(bg_pic._element)
    slide.shapes._spTree.insert(2, bg_pic._element)
    
    # Add top-right image if provided
    if top_right_img_path:
        margin_emu = margin * 9525
        img_width, img_height = 152, 94
        width_emu = img_width * 9525
        height_emu = img_height * 9525
        left_pos = p.slide_width - width_emu - margin_emu
        top_pos = margin_emu
        
        slide.shapes.add_picture(
            top_right_img_path, left_pos, top_pos, 
            width=width_emu, height=height_emu
        )
    
    add_vulnerability_content(slide, VULNERABILITIES, content_cfg={
        'x': Inches(0),
        'y': Cm(14.75),
        'width': Cm(21),
        'height': None,  # Will be calculated
        'item_height': Inches(1.1),
        'icon_size': Inches(0.6),
        'icon_margin': Inches(0.3),
        'text_left_offset': Inches(1.1),
        'text_width_ratio': 0.8,
        'items_spacing': Inches(0),
    })
    
    heading = "Cybersecurity Executive Brief​"
    text = ("Honestycar faces significant cybersecurity risks like improper input validation, weak ciphers, "
            "and insecure file upload functionality and network vulnerabilities. These weaknesses pose "
            "significant risks to data integrity, confidentiality, and business operations.​")
    add_transparent_box_with_text(
        slide, 
        heading, 
        text, 
        x=Cm(0.5),      # Adjust X position
        y=Cm(6.5),      # Adjust Y position  
        width=Cm(12), # Adjust width
        height=Cm(4)
    )
    
    add_yellow_box_custom_corners(slide)
    
    return slide